"""Office agent LLM-facing tools: document reading, CSV analysis, workspace write."""

from __future__ import annotations

import csv
import json
import os
import re
import shutil
import time
import zipfile
import xml.etree.ElementTree as ET
from html.parser import HTMLParser
from io import StringIO
from pathlib import Path

from framework.tools.base import BaseTool, ToolResult


# ---------------------------------------------------------------------------
# Path validation helper
# ---------------------------------------------------------------------------

def _get_source_root() -> str:
    return os.environ.get("OFFICE_SOURCE_ROOT", "/")


def _get_workspace_root() -> str:
    return os.environ.get("OFFICE_WORKSPACE_ROOT", "")


def _allowed_base_paths() -> list[str]:
    allowed_bases = os.environ.get("OFFICE_ALLOWED_BASE_PATHS", "")
    return [base.strip() for base in allowed_bases.split(":") if base.strip()]


def _metadata_root_for_path(path: str) -> str:
    real_path = os.path.realpath(os.path.abspath(path))
    candidates: list[str] = []
    for base in _allowed_base_paths():
        base_real = os.path.realpath(os.path.abspath(base))
        prefix = base_real.rstrip(os.sep) + os.sep
        if real_path == base_real:
            candidate = os.path.dirname(base_real) if os.path.isfile(base_real) else base_real
            if candidate:
                candidates.append(candidate)
            continue
        if real_path.startswith(prefix):
            candidates.append(base_real)

    if candidates:
        return max(candidates, key=len)
    return _get_source_root()


def _validate_path(path: str) -> tuple[str, str]:
    """Validate that path is within OFFICE_SOURCE_ROOT and OFFICE_ALLOWED_BASE_PATHS.

    Returns (normalized_path, "") on success.
    Returns ("", error_message) on failure.
    """
    source_root = os.environ.get("OFFICE_SOURCE_ROOT", "/")
    allowed_bases = os.environ.get("OFFICE_ALLOWED_BASE_PATHS", "")

    try:
        real_path = os.path.realpath(os.path.abspath(path))
        real_root = os.path.realpath(os.path.abspath(source_root))

        # If OFFICE_SOURCE_ROOT is explicitly set to "/" (root), allow all paths
        if real_root != os.sep:
            prefix = real_root.rstrip(os.sep) + os.sep
            if os.path.islink(path):
                link_target = os.path.realpath(os.path.abspath(os.readlink(path)))
                if not link_target.startswith(prefix):
                    return "", "Symlink target outside OFFICE_SOURCE_ROOT"
            # Allow exact match (accessing the source root directory itself)
            if real_path != real_root and not real_path.startswith(prefix):
                return "", f"Path {path!r} is outside OFFICE_SOURCE_ROOT"

        # Check OFFICE_ALLOWED_BASE_PATHS whitelist if set
        if allowed_bases:
            allowed_list = [bp.strip() for bp in allowed_bases.split(":") if bp.strip()]
            if allowed_list:
                in_allowed = False
                for base in allowed_list:
                    base_real = os.path.realpath(os.path.abspath(base))
                    if real_path.startswith(base_real + os.sep) or real_path == base_real:
                        in_allowed = True
                        break
                if not in_allowed:
                    return "", f"Path {path!r} is not in OFFICE_ALLOWED_BASE_PATHS. Allowed bases: {allowed_bases}"

        return real_path, ""
    except Exception as exc:
        return "", f"Path validation error: {exc}"


def _validate_workspace_path(path: str) -> tuple[str, str]:
    """Validate that a path is within OFFICE_WORKSPACE_ROOT."""
    workspace_root = _get_workspace_root()
    if not workspace_root:
        return "", "OFFICE_WORKSPACE_ROOT is not set"
    try:
        real_path = os.path.realpath(os.path.abspath(path))
        real_root = os.path.realpath(os.path.abspath(workspace_root))
        prefix = real_root.rstrip(os.sep) + os.sep
        if real_path != real_root and not real_path.startswith(prefix):
            return "", f"Path {path!r} is outside OFFICE_WORKSPACE_ROOT"
        return real_path, ""
    except Exception as exc:
        return "", f"Workspace path validation error: {exc}"


# ---------------------------------------------------------------------------
# File size limit helper
# ---------------------------------------------------------------------------

def _check_file_size(path: str) -> tuple[bool, str]:
    """Check if file exceeds size limit. Returns (ok, error_message)."""
    max_size_mb = int(os.environ.get("OFFICE_MAX_FILE_SIZE_MB", "50"))
    max_size_bytes = max_size_mb * 1024 * 1024
    try:
        size = os.path.getsize(path)
        if size > max_size_bytes:
            return False, f"File size {size / (1024*1024):.1f}MB exceeds maximum allowed size of {max_size_mb}MB. Set OFFICE_MAX_FILE_SIZE_MB to increase limit."
        return True, ""
    except OSError:
        return True, ""  # Let the file read handle the error


def _decode_text_bytes(raw: bytes) -> tuple[str, str]:
    """Decode bytes robustly without requiring optional third-party packages."""
    try:
        import chardet  # type: ignore
    except Exception:
        chardet = None

    if chardet is not None:
        try:
            detected = chardet.detect(raw)
            encoding = (detected or {}).get("encoding") or ""
            if encoding:
                return raw.decode(encoding, errors="replace"), encoding
        except Exception:
            pass

    for encoding in ("utf-8-sig", "utf-8", "utf-16", "utf-16-le", "utf-16-be", "cp1252", "latin-1"):
        try:
            return raw.decode(encoding), encoding
        except Exception:
            continue

    return raw.decode("utf-8", errors="replace"), "utf-8-replace"


def _is_number(value: object) -> bool:
    try:
        if value is None:
            return False
        text = str(value).strip().replace(",", "")
        if not text:
            return False
        float(text)
        return True
    except Exception:
        return False


def _summarize_tabular_rows(headers: list[str], rows: list[list[object]], sample_limit: int = 25) -> dict:
    """Build a compact, analysis-oriented summary for tabular data."""
    normalized_headers = [
        str(header).strip() if str(header).strip() else f"column_{idx + 1}"
        for idx, header in enumerate(headers)
    ]
    row_dicts: list[dict[str, str]] = []
    for row in rows:
        row_dicts.append({
            normalized_headers[idx]: (
                "" if idx >= len(row) or row[idx] is None else str(row[idx])
            )
            for idx in range(len(normalized_headers))
        })

    schema: list[dict[str, object]] = []
    numeric_stats: dict[str, dict[str, float | int]] = {}
    categorical_previews: dict[str, list[dict[str, object]]] = {}

    for idx, header in enumerate(normalized_headers):
        values = [
            row[idx] if idx < len(row) else ""
            for row in rows
        ]
        non_empty = [value for value in values if str(value).strip()]
        numeric_values = [float(str(value).strip().replace(",", "")) for value in non_empty if _is_number(value)]
        missing_count = len(values) - len(non_empty)
        inferred_type = "numeric" if non_empty and len(numeric_values) / len(non_empty) >= 0.8 else "text"
        schema.append({
            "name": header,
            "inferred_type": inferred_type,
            "missing_count": missing_count,
            "non_empty_count": len(non_empty),
        })

        if inferred_type == "numeric" and numeric_values:
            numeric_stats[header] = {
                "count": len(numeric_values),
                "min": min(numeric_values),
                "max": max(numeric_values),
                "avg": round(sum(numeric_values) / len(numeric_values), 2),
            }
        elif non_empty:
            counts: dict[str, int] = {}
            for value in non_empty:
                key = str(value).strip()
                counts[key] = counts.get(key, 0) + 1
            top_values = sorted(counts.items(), key=lambda item: item[1], reverse=True)[:5]
            categorical_previews[header] = [
                {"value": value, "count": count}
                for value, count in top_values
            ]

    return {
        "headers": normalized_headers,
        "total_rows": len(rows),
        "schema": schema,
        "numeric_stats": numeric_stats,
        "categorical_previews": categorical_previews,
        "sample_rows": row_dicts[:sample_limit],
    }


def _truncate_content(text: str, max_chars: int = 16000) -> tuple[str, bool]:
    """Bound tool output size so multi-document tasks stay within context/time budgets."""
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars], True


# ---------------------------------------------------------------------------
# Audit file helpers
# ---------------------------------------------------------------------------

ORGANIZED_OUTPUT_ROOT = "organized-output/files/"
VALID_CATEGORIES = {"students", "documents", "data", "code", "images", "presentations"}
WRAPPER_PREFIXES = {"grouped", "by-student", "organized", "output", "originals"}
TEXT_PREVIEW_EXTENSIONS = {
    ".txt", ".md", ".markdown",
    ".csv", ".tsv",
    ".json", ".jsonl",
    ".yaml", ".yml",
    ".rtf", ".log",
    ".ini", ".cfg", ".toml",
    ".html", ".htm", ".xml",
}
READER_TOOL_BY_EXTENSION = {
    ".txt": "read_txt",
    ".md": "read_txt",
    ".markdown": "read_txt",
    ".html": "read_txt",
    ".htm": "read_txt",
    ".xml": "read_txt",
    ".json": "read_txt",
    ".jsonl": "read_txt",
    ".yaml": "read_txt",
    ".yml": "read_txt",
    ".ini": "read_txt",
    ".cfg": "read_txt",
    ".toml": "read_txt",
    ".log": "read_txt",
    ".rtf": "read_txt",
    ".csv": "read_csv",
    ".tsv": "read_txt",
    ".pdf": "read_pdf",
    ".docx": "read_docx",
    ".docm": "read_docx",
    ".dotx": "read_docx",
    ".dotm": "read_docx",
    ".odt": "read_docx",
    ".pptx": "read_pptx",
    ".pptm": "read_pptx",
    ".potx": "read_pptx",
    ".potm": "read_pptx",
    ".ppsx": "read_pptx",
    ".ppsm": "read_pptx",
    ".odp": "read_pptx",
    ".xlsx": "read_xlsx",
    ".xlsm": "read_xlsx",
    ".xltx": "read_xlsx",
    ".xltm": "read_xlsx",
    ".xlsb": "read_xlsx",
    ".ods": "read_xlsx",
    ".xls": "read_xls",
}
IDENTITY_PREFIXES = {
    "student",
    "author",
    "writer",
    "owner",
    "employee",
    "member",
    "candidate",
    "user",
    "agent",
    "participant",
    "customer",
    "client",
    "employee name",
    "name",
    "created by",
    "prepared by",
}


def _is_wrapper_prefixed(path: str) -> bool:
    """Check if path starts with a wrapper prefix like grouped/, by-student/, etc."""
    parts = path.strip("/").split("/")
    if parts:
        return parts[0].rstrip("s") in WRAPPER_PREFIXES or parts[0] in WRAPPER_PREFIXES
    return False


def _normalize_organized_path(target_path: str, output_mode: str = "workspace") -> str:
    """Normalize organize output path to organized-output/files/ schema.

    Strips wrapper prefixes like grouped/, by-student/, etc. and ensures
    all paths are under organized-output/files/.
    """
    path = target_path.strip()
    # Strip leading slash
    if path.startswith("/"):
        path = path[1:]
    # Strip wrapper prefixes
    while _is_wrapper_prefixed(path):
        parts = path.strip("/").split("/", 1)
        path = parts[1] if len(parts) > 1 else ""
    # Add schema prefix
    if not path.startswith(ORGANIZED_OUTPUT_ROOT):
        path = ORGANIZED_OUTPUT_ROOT + path
    return path


def _is_under_organized_output(path: str) -> bool:
    """Check if path is under organized-output/files/ schema."""
    normalized = path.strip()
    if normalized.startswith("/"):
        normalized = normalized[len(ORGANIZED_OUTPUT_ROOT):]
    # Check if it's under organized-output/files/ or is a category-relative path
    if normalized.startswith(ORGANIZED_OUTPUT_ROOT) or normalized == ORGANIZED_OUTPUT_ROOT.rstrip("/"):
        return True
    # Check if it's a category-relative path (students/, documents/, data/, etc.)
    first_component = normalized.split("/")[0] if normalized else ""
    return first_component in VALID_CATEGORIES


def _categorize_extension(ext: str) -> str:
    if ext in (".pdf", ".doc", ".docx", ".docm", ".dotx", ".dotm", ".odt"):
        return "documents"
    if ext in (
        ".txt", ".md", ".markdown", ".rtf", ".html", ".htm", ".xml",
        ".json", ".jsonl", ".yaml", ".yml", ".log", ".ini", ".cfg", ".toml",
    ):
        return "text"
    if ext in (".csv", ".xlsx", ".xls", ".xlsm", ".xltx", ".xltm", ".xlsb", ".ods", ".tsv"):
        return "data"
    if ext in (".png", ".jpg", ".jpeg", ".gif", ".svg"):
        return "images"
    if ext in (".ppt", ".pptx", ".pptm", ".potx", ".potm", ".ppsx", ".ppsm", ".odp"):
        return "presentations"
    if ext in (".py", ".js", ".ts", ".java", ".cpp", ".c", ".h"):
        return "code"
    return "other"


def _suggested_reader_tool(ext: str) -> str | None:
    return READER_TOOL_BY_EXTENSION.get(ext)


def _path_year_hint(root: str, relative_path: str) -> str | None:
    candidates = list(Path(root).parts) + list(Path(relative_path).parts)
    for part in candidates:
        if re.fullmatch(r"(19|20)\d{2}", part):
            return part
    return None


def _infer_date_bucket(root: str, relative_path: str, preview: str) -> str | None:
    year_hint = _path_year_hint(root, relative_path)
    tokens = list(Path(relative_path).parts)
    tokens.extend(re.findall(r"\b\d{4}[-_/]?\d{2}(?:[-_/]?\d{2})?\b", preview))

    for token in tokens:
        cleaned = token.replace("_", "-").replace("/", "-")
        full_date = re.search(r"\b((19|20)\d{2})-(\d{2})(?:-(\d{2}))?\b", cleaned)
        if full_date:
            return f"{full_date.group(1)}-{full_date.group(3)}"
        compact_date = re.search(r"\b((19|20)\d{2})(\d{2})(\d{2})\b", token)
        if compact_date:
            return f"{compact_date.group(1)}-{compact_date.group(3)}"
        compact_month = re.search(r"\b((19|20)\d{2})(\d{2})\b", token)
        if compact_month:
            return f"{compact_month.group(1)}-{compact_month.group(3)}"
        month_day = re.fullmatch(r"(\d{2})(\d{2})", token)
        if month_day and year_hint:
            month = month_day.group(1)
            if 1 <= int(month) <= 12:
                return f"{year_hint}-{month}"
    return None


def _extract_prominent_headings(lines: list[str]) -> list[str]:
    headings: list[str] = []
    for line in lines[:20]:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith((">>>", "#", "##")):
            headings.append(stripped.lstrip("#>").strip())
            continue
        if len(stripped) <= 80 and stripped == stripped.title():
            headings.append(stripped)
        if len(headings) >= 5:
            break
    return headings


def _extract_labeled_fields(lines: list[str]) -> list[dict[str, str]]:
    fields: list[dict[str, str]] = []
    for line in lines[:25]:
        match = re.match(r"^\s*([A-Za-z][A-Za-z0-9 _/-]{1,40})\s*:\s*(.+?)\s*$", line)
        if not match:
            continue
        label = match.group(1).strip()
        value = match.group(2).strip()
        if len(value) > 120:
            continue
        fields.append({"label": label, "value": value})
        if len(fields) >= 8:
            break
    return fields


def _clean_entity_candidate(value: str) -> str:
    candidate = re.sub(r"^[>\-#\s]+", "", value).strip()
    candidate = re.sub(r"\s+", " ", candidate)
    lowered = candidate.lower()
    for prefix in sorted(IDENTITY_PREFIXES, key=len, reverse=True):
        if lowered.startswith(prefix + " "):
            return candidate[len(prefix):].strip(" :-")
    return candidate.strip(" :-")


def _looks_like_person_name(value: str) -> bool:
    tokens = re.findall(r"[A-Za-z][A-Za-z'._-]*", value)
    if not tokens or len(tokens) > 4:
        return False
    return all(token[:1].isupper() for token in tokens if token)


def _extract_primary_entity(lines: list[str], headings: list[str], labeled_fields: list[dict[str, str]]) -> tuple[str | None, str | None, str]:
    for line in lines[:20]:
        stripped = line.strip()
        for prefix in sorted(IDENTITY_PREFIXES, key=len, reverse=True):
            marker = prefix.title()
            if stripped.lower().startswith((">>> " + prefix + " ", prefix + " ")):
                candidate = _clean_entity_candidate(
                    stripped.split(">>>", 1)[-1].strip() if stripped.startswith(">>>") else stripped
                )
                if candidate.lower().startswith(prefix + " "):
                    candidate = candidate[len(prefix):].strip(" :-")
                if candidate:
                    return candidate, "explicit_heading", "high"
    for field in labeled_fields:
        candidate = _clean_entity_candidate(field["value"])
        label = field["label"].strip().lower()
        if label in IDENTITY_PREFIXES and candidate:
            return candidate, "labeled_field", "high"
    return None, None, "none"


def _read_text_preview(path: str, max_chars: int = 1200) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext not in TEXT_PREVIEW_EXTENSIONS:
        return ""
    try:
        with open(path, "rb") as fh:
            raw = fh.read(max_chars * 2)
        text, _ = _decode_text_bytes(raw)
        return text[:max_chars]
    except Exception:
        return ""


def _extract_markup_text(raw_text: str, ext: str) -> tuple[str, str]:
    """Convert markup-oriented formats into readable text."""
    if ext in {".html", ".htm"}:
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(raw_text, "lxml")
            return soup.get_text("\n"), "html-bs4"
        except Exception:
            try:
                class _TextExtractor(HTMLParser):
                    def __init__(self) -> None:
                        super().__init__()
                        self.chunks: list[str] = []

                    def handle_data(self, data: str) -> None:
                        if data.strip():
                            self.chunks.append(data.strip())

                parser = _TextExtractor()
                parser.feed(raw_text)
                return "\n".join(parser.chunks), "html-stdlib"
            except Exception:
                return raw_text, "html-raw"
    if ext == ".xml":
        try:
            root = ET.fromstring(raw_text)
            return "\n".join(text.strip() for text in root.itertext() if text.strip()), "xml-etree"
        except Exception:
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(raw_text, "xml")
                return soup.get_text("\n"), "xml-bs4"
            except Exception:
                return raw_text, "xml-raw"
    if ext == ".json":
        try:
            parsed = json.loads(raw_text)
            return json.dumps(parsed, ensure_ascii=False, indent=2), "json-pretty"
        except Exception:
            return raw_text, "json-raw"
    if ext == ".jsonl":
        lines = []
        for raw_line in raw_text.splitlines():
            raw_line = raw_line.strip()
            if not raw_line:
                continue
            try:
                lines.append(json.dumps(json.loads(raw_line), ensure_ascii=False, indent=2))
            except Exception:
                lines.append(raw_line)
        return "\n\n".join(lines), "jsonl-pretty"
    if ext in {".yaml", ".yml"}:
        try:
            import yaml
            parsed = yaml.safe_load(raw_text)
            return json.dumps(parsed, ensure_ascii=False, indent=2), "yaml-json"
        except Exception:
            return raw_text, "yaml-raw"
    if ext == ".rtf":
        try:
            from striprtf.striprtf import rtf_to_text
            return rtf_to_text(raw_text), "rtf-striprtf"
        except Exception:
            return raw_text, "rtf-raw"
    return raw_text, "plain-text"


def _extract_odf_text_nodes(content_xml: bytes, tags: tuple[str, ...]) -> list[str]:
    root = ET.fromstring(content_xml)
    paragraphs: list[str] = []
    for element in root.iter():
        tag_name = element.tag.rsplit("}", 1)[-1]
        if tag_name not in tags:
            continue
        joined = "".join(text.strip() for text in element.itertext() if text and text.strip()).strip()
        if joined:
            paragraphs.append(joined)
    return paragraphs


def _extract_docx_like_text(path: str) -> tuple[list[str], str]:
    """Extract paragraph text from Word-like documents without OCR."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".odt":
        with zipfile.ZipFile(path) as archive:
            xml_bytes = archive.read("content.xml")
        return _extract_odf_text_nodes(xml_bytes, ("h", "p")), "odt-zip-xml"
    try:
        import docx
        doc = docx.Document(path)
        return [p.text for p in doc.paragraphs if p.text.strip()], "python-docx"
    except Exception:
        with zipfile.ZipFile(path) as archive:
            xml_bytes = archive.read("word/document.xml")
        root = ET.fromstring(xml_bytes)
        namespaces = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        paragraphs = []
        for para in root.findall(".//w:p", namespaces):
            texts = [node.text or "" for node in para.findall(".//w:t", namespaces)]
            joined = "".join(texts).strip()
            if joined:
                paragraphs.append(joined)
        return paragraphs, "zip-xml"


def _extract_presentation_like_text(path: str) -> tuple[list[str], str, int]:
    """Extract slide-like text from PowerPoint-like documents without OCR."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".odp":
        with zipfile.ZipFile(path) as archive:
            xml_bytes = archive.read("content.xml")
        slides = _extract_odf_text_nodes(xml_bytes, ("h", "p"))
        return slides, "odp-zip-xml", len(slides)

    import pptx

    prs = pptx.Presentation(path)
    slides: list[str] = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            slides.append(f"[Slide {i + 1}]\n" + "\n".join(slide_text))
    return slides, "python-pptx", len(prs.slides)


def _extract_ods_rows(path: str) -> tuple[dict[str, dict[str, object]], list[str]]:
    with zipfile.ZipFile(path) as archive:
        xml_bytes = archive.read("content.xml")
    root = ET.fromstring(xml_bytes)
    sheets: dict[str, dict[str, object]] = {}
    sheet_names: list[str] = []
    table_ns = "urn:oasis:names:tc:opendocument:xmlns:table:1.0"
    text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    for table in root.findall(f".//{{{table_ns}}}table"):
        sheet_name = table.attrib.get(f"{{{table_ns}}}name", f"Sheet{len(sheet_names) + 1}")
        sheet_names.append(sheet_name)
        rows: list[list[str]] = []
        for row_el in table.findall(f"./{{{table_ns}}}table-row"):
            repeat_rows = int(row_el.attrib.get(f"{{{table_ns}}}number-rows-repeated", "1"))
            row_values: list[str] = []
            for cell in row_el.findall(f"./{{{table_ns}}}table-cell"):
                repeat_cells = int(cell.attrib.get(f"{{{table_ns}}}number-columns-repeated", "1"))
                cell_text = "\n".join(
                    text.strip()
                    for text in cell.itertext()
                    if text and text.strip()
                )
                row_values.extend([cell_text] * repeat_cells)
            if row_values:
                for _ in range(min(repeat_rows, 1)):
                    rows.append(row_values.copy())
        if not rows:
            sheets[sheet_name] = {"headers": [], "total_rows": 0, "sample_rows": []}
            continue
        sheets[sheet_name] = _summarize_tabular_rows(rows[0], rows[1:])
    return sheets, sheet_names


def _safe_path_segment(value: str) -> str:
    segment = re.sub(r"[^A-Za-z0-9._-]+", "_", value.strip())
    return segment.strip("_") or "unknown"


def _build_file_metadata(root: str, full_path: str) -> dict[str, object]:
    rel_path = os.path.relpath(full_path, root)
    ext = os.path.splitext(full_path)[1].lower()
    category = _categorize_extension(ext)
    preview = _read_text_preview(full_path)
    lines = preview.splitlines()
    headings = _extract_prominent_headings(lines)
    labeled_fields = _extract_labeled_fields(lines)
    primary_entity, primary_entity_source, primary_entity_confidence = _extract_primary_entity(
        lines,
        headings,
        labeled_fields,
    )
    inferred_date_bucket = _infer_date_bucket(root, rel_path, preview)
    relative_stem = rel_path.replace(os.sep, "-")
    suggested_destination = None
    if primary_entity and inferred_date_bucket:
        suggested_destination = (
            f"{_safe_path_segment(primary_entity)}/"
            f"{inferred_date_bucket}/"
            f"{relative_stem}"
        )
    return {
        "relative_path": rel_path,
        "name": os.path.basename(full_path),
        "ext": ext,
        "size": os.path.getsize(full_path),
        "category": category,
        "parent_dirs": list(Path(rel_path).parts[:-1]),
        "suggested_reader_tool": _suggested_reader_tool(ext),
        "inferred_date_bucket": inferred_date_bucket,
        "primary_entity": primary_entity,
        "primary_entity_source": primary_entity_source,
        "primary_entity_confidence": primary_entity_confidence,
        "prominent_headings": headings[:2],
        "labeled_fields": labeled_fields[:2],
        "suggested_destination": suggested_destination,
    }


def collect_organize_file_inventory(root: str) -> tuple[list[dict[str, object]], dict[str, list[str]], int]:
    """Return recursive file inventory and category groups for organize tasks."""
    inventory: list[dict[str, object]] = []
    groups: dict[str, list[str]] = {}
    total_dirs = 0
    for walk_root, dirs, files in os.walk(root):
        dirs[:] = sorted(d for d in dirs if not d.startswith("."))
        total_dirs += len(dirs)
        for name in sorted(files):
            if name.startswith("."):
                continue
            full_path = os.path.join(walk_root, name)
            item = _build_file_metadata(root, full_path)
            inventory.append(item)
            groups.setdefault(str(item["category"]), []).append(str(item["relative_path"]))
    inventory.sort(key=lambda item: str(item["relative_path"]))
    return inventory, groups, total_dirs


def _get_audit_dir() -> str:
    """Return the audit directory path for organize operations."""
    audit_dir = os.environ.get("OFFICE_AUDIT_DIR", "")
    if audit_dir:
        return audit_dir
    workspace_root = os.environ.get("OFFICE_WORKSPACE_ROOT", "")
    if workspace_root:
        return workspace_root
    artifact_root = os.environ.get("ARTIFACT_ROOT", "/tmp")
    return os.path.join(artifact_root, "office", "audit")


def _backup_existing_file(path: str) -> str | None:
    """Backup existing file with timestamp suffix. Returns backup path or None."""
    backup_enabled = os.environ.get("OFFICE_BACKUP_ENABLED", "true").lower()
    if backup_enabled not in ("true", "1", "yes"):
        return None
    if not os.path.exists(path):
        return None
    timestamp = time.strftime("%Y%m%d-%H%M%S")
    backup_path = f"{path}.{timestamp}.bak"
    try:
        os.rename(path, backup_path)
        return backup_path
    except OSError:
        return None


def _write_operations_plan(steps: list[dict], source_paths: list[str], output_mode: str, task_id: str = "") -> str:
    """Write operations-plan.json before any file operation begins."""
    audit_dir = _get_audit_dir()
    os.makedirs(audit_dir, exist_ok=True)
    plan_path = os.path.join(audit_dir, "operations-plan.json")
    plan = {
        "version": "1.0",
        "timestamp": time.strftime("%Y-%m-%dT%H:%M:%S%z"),
        "capability": "organize",
        "task_id": task_id,
        "source_paths": source_paths,
        "output_mode": output_mode,
        "steps": [{"step_id": i + 1, "status": "pending", **s} for i, s in enumerate(steps)],
    }
    with open(plan_path, "w", encoding="utf-8") as f:
        json.dump(plan, f, ensure_ascii=False, indent=2)
    return plan_path


def _update_step_status(step_id: int, status: str, plan_path: str | None = None) -> None:
    """Update a step's status in operations-plan.json."""
    if not plan_path or not os.path.exists(plan_path):
        return
    try:
        with open(plan_path, encoding="utf-8") as f:
            plan = json.load(f)
        for step in plan.get("steps", []):
            if step.get("step_id") == step_id:
                step["status"] = status
                break
        with open(plan_path, "w", encoding="utf-8") as f:
            json.dump(plan, f, ensure_ascii=False, indent=2)
    except Exception:
        pass


def _append_command_log(action: str, params: dict, step_id: int | None = None) -> None:
    """Append a timestamped operation entry to command-log.txt."""
    audit_dir = _get_audit_dir()
    os.makedirs(audit_dir, exist_ok=True)
    log_path = os.path.join(audit_dir, "command-log.txt")
    entry = f"[{time.strftime('%Y-%m-%dT%H:%M:%S%z')}] STEP {step_id or '?'}: {action} {json.dumps(params, ensure_ascii=False)}\n"
    try:
        with open(log_path, "a", encoding="utf-8") as f:
            f.write(entry)
    except Exception:
        pass


def _write_stage_summary(stage: str, completed_steps: list, pending_steps: list, warnings: list, errors: list) -> str:
    """Write stage-summary.json after each major stage."""
    audit_dir = _get_audit_dir()
    os.makedirs(audit_dir, exist_ok=True)
    summary_path = os.path.join(audit_dir, "stage-summary.json")
    summary = {
        "timestamp": time.strftime("%Y-%m-%dT%H:%M:%S%z"),
        "stage": stage,
        "completed_steps": completed_steps,
        "pending_steps": pending_steps,
        "warnings": warnings,
        "errors": errors,
    }
    with open(summary_path, "w", encoding="utf-8") as f:
        json.dump(summary, f, ensure_ascii=False, indent=2)
    return summary_path


def _write_manifest_entry(action: str, target: str, status: str, task_id: str = "") -> str | None:
    """Write/update .office-agent-manifest.json incrementally for in-place mode."""
    audit_dir = os.environ.get("OFFICE_SOURCE_ROOT", "")
    if not audit_dir:
        return None
    manifest_path = os.path.join(audit_dir, ".office-agent-manifest.json")
    manifest = {"version": "1.0", "task_id": task_id, "output_mode": "inplace", "operations": []}
    if os.path.exists(manifest_path):
        try:
            with open(manifest_path, encoding="utf-8") as f:
                manifest = json.load(f)
        except Exception:
            pass
    manifest["operations"].append({"action": action, "target": target, "status": status, "timestamp": time.strftime("%Y-%m-%dT%H:%M:%S%z")})
    manifest["rollback_possible"] = True
    try:
        with open(manifest_path, "w", encoding="utf-8") as f:
            json.dump(manifest, f, ensure_ascii=False, indent=2)
        return manifest_path
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Directory resource pre-check
# ---------------------------------------------------------------------------

def _scan_directory_resources(path: str) -> tuple[int, int]:
    """Scan directory tree and return (file_count, total_bytes).

    Skips files > 50MB and hidden files/directories.
    """
    file_count = 0
    total_bytes = 0
    max_file_size = 50 * 1024 * 1024  # 50MB
    try:
        for root, dirs, files in os.walk(path):
            # Skip hidden directories
            dirs[:] = [d for d in dirs if not d.startswith(".")]
            for fname in files:
                if fname.startswith("."):
                    continue
                fpath = os.path.join(root, fname)
                try:
                    size = os.path.getsize(fpath)
                    if size > max_file_size:
                        continue  # Skip oversized files
                    file_count += 1
                    total_bytes += size
                except OSError:
                    continue
    except OSError:
        pass
    return file_count, total_bytes


def _check_directory_limits(path: str) -> dict | None:
    """Check if directory exceeds resource limits.

    Returns error dict if exceeded, None if OK.
    """
    max_files = int(os.environ.get("OFFICE_MAX_TOTAL_FILES", "1000"))
    max_bytes = int(os.environ.get("OFFICE_MAX_TOTAL_BYTES", str(500 * 1024 * 1024)))
    file_count, total_bytes = _scan_directory_resources(path)
    if file_count > max_files:
        return {
            "error": "Resource limit exceeded",
            "pre_check_report": {
                "total_files": file_count,
                "max_files": max_files,
                "total_bytes": total_bytes,
                "max_bytes": max_bytes,
                "limit_type": "file_count",
            },
        }
    if total_bytes > max_bytes:
        return {
            "error": "Resource limit exceeded",
            "pre_check_report": {
                "total_files": file_count,
                "max_files": max_files,
                "total_bytes": total_bytes,
                "max_bytes": max_bytes,
                "limit_type": "total_bytes",
            },
        }
    return None


# ---------------------------------------------------------------------------
# Read PDF Tool
# ---------------------------------------------------------------------------

class ReadPdfTool(BaseTool):
    name = "read_pdf"
    description = "Read a PDF file and return bounded text content with extraction metadata for summarization. If no embedded text is extractable without OCR, report that clearly."
    parameters_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Absolute path to the PDF file"},
        },
        "required": ["path"],
    }

    def execute_sync(self, path: str = "") -> ToolResult:
        normalized, err = _validate_path(path)
        if err:
            return ToolResult(output="", error=f"read_pdf: {err}")
        if not os.path.isfile(normalized):
            return ToolResult(output="", error=f"read_pdf: file not found: {path}")
        if not normalized.lower().endswith(".pdf"):
            return ToolResult(output="", error=f"read_pdf: not a PDF file: {path}")
        ok, size_err = _check_file_size(normalized)
        if not ok:
            return ToolResult(output="", error=f"read_pdf: {size_err}")
        try:
            import pdfplumber
            with pdfplumber.open(normalized) as pdf:
                total_pages = len(pdf.pages)
                pages = []
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append(f"[Page {page.page_number}]\n{text}")
                content = "\n\n".join(pages)
            extraction_method = "pdfplumber"
            if not content.strip():
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(normalized)
                    alt_pages = []
                    for idx, page in enumerate(reader.pages):
                        text = page.extract_text() or ""
                        if text.strip():
                            alt_pages.append(f"[Page {idx + 1}]\n{text}")
                    if alt_pages:
                        content = "\n\n".join(alt_pages)
                        extraction_method = "pypdf"
                except Exception:
                    pass
            content, truncated = _truncate_content(content)
            size_kb = os.path.getsize(normalized) // 1024
            return ToolResult(output=json.dumps({
                "content": content,
                "path": normalized,
                "pages_with_text": len(pages) if extraction_method == "pdfplumber" else len([block for block in content.split("\n\n") if block.strip()]),
                "total_pages": total_pages,
                "size_kb": size_kb,
                "extractable_text": bool(content.strip()),
                "extraction_method": extraction_method if content.strip() else "none",
                "truncated": truncated,
            }))
        except Exception as exc:
            return ToolResult(output="", error=f"read_pdf: failed to read: {exc}")


# ---------------------------------------------------------------------------
# Read DOCX Tool
# ---------------------------------------------------------------------------

class ReadDocxTool(BaseTool):
    name = "read_docx"
    description = "Read Word-like text documents (.docx/.docm/.dotx/.dotm/.odt) and return bounded text content including paragraph metadata. Uses non-OCR fallbacks when optional libraries are unavailable."
    parameters_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Absolute path to the Word OpenXML file"},
        },
        "required": ["path"],
    }

    def execute_sync(self, path: str = "") -> ToolResult:
        normalized, err = _validate_path(path)
        if err:
            return ToolResult(output="", error=f"read_docx: {err}")
        if not os.path.isfile(normalized):
            return ToolResult(output="", error=f"read_docx: file not found: {path}")
        if normalized.lower().endswith(".doc"):
            return ToolResult(output="", error="read_docx: .doc format is not supported. Please convert to .docx first.")
        if not normalized.lower().endswith((".docx", ".docm", ".dotx", ".dotm", ".odt")):
            return ToolResult(output="", error=f"read_docx: not a supported Word-like file: {path}")
        ok, size_err = _check_file_size(normalized)
        if not ok:
            return ToolResult(output="", error=f"read_docx: {size_err}")
        try:
            paragraphs, extraction_method = _extract_docx_like_text(normalized)
            content = "\n\n".join(paragraphs)
            content, truncated = _truncate_content(content)
            return ToolResult(output=json.dumps({
                "content": content,
                "path": normalized,
                "paragraphs": len(paragraphs),
                "extraction_method": extraction_method,
                "truncated": truncated,
            }))
        except Exception as exc:
            return ToolResult(output="", error=f"read_docx: failed to read: {exc}")


# ---------------------------------------------------------------------------
# Read PPTX Tool
# ---------------------------------------------------------------------------

class ReadPptxTool(BaseTool):
    name = "read_pptx"
    description = "Read PowerPoint-like files (.pptx/.pptm/.potx/.potm/.ppsx/.ppsm/.odp) and return bounded slide text for summarization."
    parameters_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Absolute path to the PPTX file"},
        },
        "required": ["path"],
    }

    def execute_sync(self, path: str = "") -> ToolResult:
        normalized, err = _validate_path(path)
        if err:
            return ToolResult(output="", error=f"read_pptx: {err}")
        if not os.path.isfile(normalized):
            return ToolResult(output="", error=f"read_pptx: file not found: {path}")
        if normalized.lower().endswith(".ppt"):
            return ToolResult(output="", error="read_pptx: legacy .ppt files are not supported. Please convert to .pptx first.")
        if not normalized.lower().endswith((".pptx", ".pptm", ".potx", ".potm", ".ppsx", ".ppsm", ".odp")):
            return ToolResult(output="", error="read_pptx: not a supported presentation file.")
        ok, size_err = _check_file_size(normalized)
        if not ok:
            return ToolResult(output="", error=f"read_pptx: {size_err}")
        try:
            slides, extraction_method, total_slides = _extract_presentation_like_text(normalized)
            content = "\n\n".join(slides)
            content, truncated = _truncate_content(content)
            return ToolResult(output=json.dumps({
                "content": content,
                "path": normalized,
                "slides": len(slides),
                "total_slides": total_slides,
                "extraction_method": extraction_method,
                "truncated": truncated,
            }))
        except Exception as exc:
            return ToolResult(output="", error=f"read_pptx: failed to read: {exc}")


# ---------------------------------------------------------------------------
# Read TXT Tool
# ---------------------------------------------------------------------------

class ReadTxtTool(BaseTool):
    name = "read_txt"
    description = "Read text-like files such as TXT, Markdown, HTML, and XML, returning bounded readable text content."
    parameters_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Absolute path to the TXT file"},
        },
        "required": ["path"],
    }

    def execute_sync(self, path: str = "") -> ToolResult:
        normalized, err = _validate_path(path)
        if err:
            return ToolResult(output="", error=f"read_txt: {err}")
        if not os.path.isfile(normalized):
            return ToolResult(output="", error=f"read_txt: file not found: {path}")
        ok, size_err = _check_file_size(normalized)
        if not ok:
            return ToolResult(output="", error=f"read_txt: {size_err}")
        try:
            with open(normalized, "rb") as fh:
                raw = fh.read()
            content, encoding = _decode_text_bytes(raw)
            ext = os.path.splitext(normalized)[1].lower()
            content, extraction_method = _extract_markup_text(content, ext)
            content, truncated = _truncate_content(content)
            return ToolResult(output=json.dumps({
                "content": content,
                "path": normalized,
                "chars": len(content),
                "encoding": encoding,
                "extraction_method": extraction_method,
                "truncated": truncated,
            }))
        except Exception as exc:
            return ToolResult(output="", error=f"read_txt: failed to read: {exc}")


# ---------------------------------------------------------------------------
# Read CSV Tool
# ---------------------------------------------------------------------------

class ReadCsvTool(BaseTool):
    name = "read_csv"
    description = "Read a CSV file and return a compact schema-first summary with sample rows and statistics."
    parameters_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Absolute path to the CSV file"},
        },
        "required": ["path"],
    }

    def execute_sync(self, path: str = "") -> ToolResult:
        normalized, err = _validate_path(path)
        if err:
            return ToolResult(output="", error=f"read_csv: {err}")
        if not os.path.isfile(normalized):
            return ToolResult(output="", error=f"read_csv: file not found: {path}")
        if not normalized.lower().endswith(".csv"):
            return ToolResult(output="", error=f"read_csv: not a CSV file: {path}")
        ok, size_err = _check_file_size(normalized)
        if not ok:
            return ToolResult(output="", error=f"read_csv: {size_err}")
        try:
            with open(normalized, "rb") as fh:
                raw = fh.read()
            text, encoding = _decode_text_bytes(raw)
            sample = text[:8192]
            try:
                dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|")
            except Exception:
                dialect = csv.excel
            reader = csv.reader(StringIO(text), dialect=dialect)
            rows = list(reader)
            if not rows:
                return ToolResult(output="", error="read_csv: CSV file is empty")
            headers = rows[0]
            data_rows = rows[1:]
            summary = _summarize_tabular_rows(headers, data_rows)
            summary["encoding"] = encoding
            return ToolResult(output=json.dumps(summary))
        except Exception as exc:
            return ToolResult(output="", error=f"read_csv: failed to read: {exc}")


# ---------------------------------------------------------------------------
# Read XLSX Tool
# ---------------------------------------------------------------------------

class ReadXlsxTool(BaseTool):
    name = "read_xlsx"
    description = "Read spreadsheet files (.xlsx/.xlsm/.xltx/.xltm/.xlsb/.ods) and return a compact summary for each sheet."
    parameters_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Absolute path to the XLSX file"},
        },
        "required": ["path"],
    }

    def execute_sync(self, path: str = "") -> ToolResult:
        normalized, err = _validate_path(path)
        if err:
            return ToolResult(output="", error=f"read_xlsx: {err}")
        if not os.path.isfile(normalized):
            return ToolResult(output="", error=f"read_xlsx: file not found: {path}")
        if not normalized.lower().endswith((".xlsx", ".xlsm", ".xltx", ".xltm", ".xlsb", ".ods")):
            return ToolResult(output="", error=f"read_xlsx: not a supported spreadsheet file: {path}")
        ok, size_err = _check_file_size(normalized)
        if not ok:
            return ToolResult(output="", error=f"read_xlsx: {size_err}")
        try:
            ext = os.path.splitext(normalized)[1].lower()
            extraction_method = ""
            if ext in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
                import openpyxl

                wb = openpyxl.load_workbook(normalized, read_only=True, data_only=True)
                sheets = {}
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    rows = [list(row) for row in ws.iter_rows(values_only=True)]
                    if not rows:
                        sheets[sheet_name] = {"headers": [], "total_rows": 0, "sample_rows": []}
                        continue
                    sheets[sheet_name] = _summarize_tabular_rows(rows[0], rows[1:])
                sheet_names = wb.sheetnames
                extraction_method = "openpyxl"
            elif ext == ".xlsb":
                from pyxlsb import open_workbook

                sheets = {}
                sheet_names = []
                with open_workbook(normalized) as wb:
                    for sheet_name in wb.sheets:
                        sheet_names.append(sheet_name)
                        with wb.get_sheet(sheet_name) as ws:
                            rows = [[cell.v for cell in row] for row in ws.rows()]
                        if not rows:
                            sheets[sheet_name] = {"headers": [], "total_rows": 0, "sample_rows": []}
                            continue
                        sheets[sheet_name] = _summarize_tabular_rows(rows[0], rows[1:])
                extraction_method = "pyxlsb"
            else:
                sheets, sheet_names = _extract_ods_rows(normalized)
                extraction_method = "ods-zip-xml"
            return ToolResult(output=json.dumps({
                "path": normalized,
                "sheets": sheets,
                "sheet_names": sheet_names,
                "extraction_method": extraction_method,
            }))
        except Exception as exc:
            return ToolResult(output="", error=f"read_xlsx: failed to read: {exc}")


# ---------------------------------------------------------------------------
# Read XLS Tool
# ---------------------------------------------------------------------------

class ReadXlsTool(BaseTool):
    name = "read_xls"
    description = "Read a legacy Excel XLS file and return a compact summary for each sheet."
    parameters_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Absolute path to the XLS file"},
        },
        "required": ["path"],
    }

    def execute_sync(self, path: str = "") -> ToolResult:
        normalized, err = _validate_path(path)
        if err:
            return ToolResult(output="", error=f"read_xls: {err}")
        if not os.path.isfile(normalized):
            return ToolResult(output="", error=f"read_xls: file not found: {path}")
        if not normalized.lower().endswith(".xls"):
            return ToolResult(output="", error=f"read_xls: not an XLS file: {path}")
        ok, size_err = _check_file_size(normalized)
        if not ok:
            return ToolResult(output="", error=f"read_xls: {size_err}")
        try:
            import xlrd
            wb = xlrd.open_workbook(normalized)
            sheets = {}
            for i in range(wb.nsheets):
                ws = wb.sheet_by_index(i)
                rows = [ws.row_values(row_idx) for row_idx in range(ws.nrows)]
                if not rows:
                    sheets[ws.name] = {"headers": [], "total_rows": 0, "sample_rows": []}
                    continue
                sheets[ws.name] = _summarize_tabular_rows(rows[0], rows[1:])
            return ToolResult(output=json.dumps({
                "path": normalized,
                "sheets": sheets,
                "sheet_names": wb.sheet_names(),
                "warning": "XLS format has limited support. Convert to XLSX for better results.",
            }))
        except ImportError:
            return ToolResult(output="", error="read_xls: xlrd library not installed. Convert the XLS file to XLSX format.")
        except Exception as exc:
            return ToolResult(output="", error=f"read_xls: failed to read (best-effort): {exc}")


# ---------------------------------------------------------------------------
# List Directory Tool
# ---------------------------------------------------------------------------

class ListDirectoryTool(BaseTool):
    name = "list_directory"
    description = "List the files and subdirectories in a given path, with file sizes and types."
    parameters_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Absolute path to the directory"},
        },
        "required": ["path"],
    }

    def execute_sync(self, path: str = "") -> ToolResult:
        normalized, err = _validate_path(path)
        if err:
            return ToolResult(output="", error=f"list_directory: {err}")
        if not os.path.isdir(normalized):
            return ToolResult(output="", error=f"list_directory: not a directory: {path}")
        try:
            entries = []
            for name in os.listdir(normalized):
                full = os.path.join(normalized, name)
                try:
                    stat = os.stat(full)
                    entries.append({
                        "name": name,
                        "type": "dir" if os.path.isdir(full) else "file",
                        "size": stat.st_size if os.path.isfile(full) else 0,
                    })
                except OSError:
                    pass
            return ToolResult(output=json.dumps({"files": entries}))
        except Exception as exc:
            return ToolResult(output="", error=f"list_directory: failed: {exc}")


# ---------------------------------------------------------------------------
# Write Workspace Tool
# ---------------------------------------------------------------------------

class WriteWorkspaceTool(BaseTool):
    name = "write_workspace"
    description = "Write content to a file in the workspace artifacts folder. Use this for workspace output mode."
    parameters_schema = {
        "type": "object",
        "properties": {
            "filename": {"type": "string", "description": "Output filename (e.g., summary.md, analysis.md)"},
            "content": {"type": "string", "description": "File content to write"},
        },
        "required": ["filename", "content"],
    }

    def execute_sync(self, filename: str = "", content: str = "") -> ToolResult:
        workspace_root = os.environ.get(
            "OFFICE_WORKSPACE_ROOT",
            os.path.join(os.environ.get("ARTIFACT_ROOT", "/tmp"), "office")
        )
        try:
            os.makedirs(workspace_root, exist_ok=True)
            safe_name = os.path.basename(filename)
            if not safe_name or safe_name != filename:
                return ToolResult(output="", error=f"write_workspace: invalid filename {filename!r}")
            out_path = os.path.join(workspace_root, safe_name)
            backup_path = _backup_existing_file(out_path)
            with open(out_path, "w", encoding="utf-8") as fh:
                fh.write(content)
            result = {"path": out_path, "bytes": len(content.encode("utf-8"))}
            if backup_path:
                result["backup_path"] = backup_path
            return ToolResult(output=json.dumps(result))
        except Exception as exc:
            return ToolResult(output="", error=f"write_workspace: failed: {exc}")


# ---------------------------------------------------------------------------
# Write File Tool (inplace mode — requires write grant)
# ---------------------------------------------------------------------------

class WriteFileTool(BaseTool):
    name = "write_file"
    description = "Write content to a file at the specified path. Only available when inplace mode is enabled with write grant."
    parameters_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Absolute path to write to (must be under OFFICE_SOURCE_ROOT)"},
            "content": {"type": "string", "description": "File content to write"},
        },
        "required": ["path", "content"],
    }

    def execute_sync(self, path: str = "", content: str = "") -> ToolResult:
        allow_inplace = os.environ.get("OFFICE_ALLOW_INPLACE_WRITES", "false").lower()
        if allow_inplace not in ("true", "1", "yes"):
            return ToolResult(output="", error="write_file: inplace mode not enabled. Set OFFICE_ALLOW_INPLACE_WRITES=true")
        normalized, err = _validate_path(path)
        if err:
            return ToolResult(output="", error=f"write_file: {err}")
        backup_path = _backup_existing_file(normalized)
        try:
            os.makedirs(os.path.dirname(normalized), exist_ok=True)
            with open(normalized, "w", encoding="utf-8") as fh:
                fh.write(content)
            result = {"path": normalized, "bytes": len(content.encode("utf-8"))}
            if backup_path:
                result["backup_path"] = backup_path
            return ToolResult(output=json.dumps(result))
        except Exception as exc:
            return ToolResult(output="", error=f"write_file: failed: {exc}")


class OrganizeFolderTool(BaseTool):
    name = "organize_folder"
    description = """Survey a folder recursively and return per-file metadata for organizing.
Returns path, recursive file inventory, category groups, and extracted signals such as headings,
date buckets, and suggested reader tools. Use this before planning organize_move_file operations."""
    parameters_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Absolute path to the folder to organize"},
        },
        "required": ["path"],
    }

    def execute_sync(self, path: str = "") -> ToolResult:
        normalized, err = _validate_path(path)
        if err:
            return ToolResult(output="", error=f"organize_folder: {err}")
        if not os.path.isdir(normalized):
            return ToolResult(output="", error=f"organize_folder: not a directory: {path}")

        try:
            files, groups, total_dirs = collect_organize_file_inventory(normalized)
            entity_counts: dict[str, int] = {}
            date_bucket_counts: dict[str, int] = {}
            for item in files:
                entity = str(item.get("primary_entity") or "").strip()
                if entity:
                    entity_counts[entity] = entity_counts.get(entity, 0) + 1
                date_bucket = str(item.get("inferred_date_bucket") or "").strip()
                if date_bucket:
                    date_bucket_counts[date_bucket] = date_bucket_counts.get(date_bucket, 0) + 1
            return ToolResult(output=json.dumps({
                "path": normalized,
                "groups": groups,
                "files": files,
                "total_files": len(files),
                "total_dirs": total_dirs,
                "entity_counts": entity_counts,
                "date_bucket_counts": date_bucket_counts,
                "errors": [],
            }))
        except Exception as exc:
            return ToolResult(output="", error=f"organize_folder: failed: {exc}")


class OrganizeMoveFileTool(BaseTool):
    name = "organize_move_file"
    description = """Execute a planned file move as part of folder organization.
Before executing, validates: (1) action is in whitelist (mkdir, copy_file, write_text),
(2) all paths are within OFFICE_SOURCE_ROOT, (3) plan is written to operations-plan.json first.
Use organize_folder tool first to survey the folder, then organize_execute_plan to run."""
    parameters_schema = {
        "type": "object",
        "properties": {
            "action": {"type": "string", "description": "Action type: mkdir, copy_file, or write_text"},
            "src": {"type": "string", "description": "Source path (for copy_file)"},
            "dst": {"type": "string", "description": "Destination path"},
            "content": {"type": "string", "description": "File content (for write_text action)"},
        },
        "required": ["action", "dst"],
    }

    ALLOWED_ACTIONS = {"mkdir", "copy_file", "write_text"}

    def execute_sync(self, action: str = "", src: str = "", dst: str = "", content: str = "") -> ToolResult:
        allow_inplace = os.environ.get("OFFICE_ALLOW_INPLACE_WRITES", "false").lower() in ("true", "1", "yes")
        workspace_root = _get_workspace_root()
        source_root = _get_source_root()
        output_mode = os.environ.get("OFFICE_OUTPUT_MODE", "workspace").lower()
        if output_mode != "inplace" and not workspace_root and allow_inplace and source_root:
            output_mode = "inplace"

        # Validate action is in whitelist
        if action not in self.ALLOWED_ACTIONS:
            return ToolResult(output="", error=f"organize_move_file: action {action!r} not allowed. Allowed: {self.ALLOWED_ACTIONS}")

        if not os.path.isabs(dst) and _is_wrapper_prefixed(dst):
            return ToolResult(output="", error=f"organize_move_file: destination {dst!r} is outside the organized-output/files/ schema (wrapper prefix not allowed)")

        # Validate destination path
        if output_mode == "inplace":
            if not allow_inplace:
                return ToolResult(output="", error="organize_move_file: inplace writes not enabled. Set OFFICE_ALLOW_INPLACE_WRITES=true")
            raw_dst = dst
            if not os.path.isabs(raw_dst):
                raw_dst = os.path.join(source_root, _normalize_organized_path(raw_dst))
            dst_normalized, err = _validate_path(raw_dst)
            if err:
                return ToolResult(output="", error=f"organize_move_file: destination {err}")
        else:
            if not workspace_root:
                return ToolResult(output="", error="organize_move_file: OFFICE_WORKSPACE_ROOT is not set for workspace mode")
            if os.path.isabs(dst):
                if not dst.startswith(workspace_root):
                    return ToolResult(output="", error=f"organize_move_file: destination {dst!r} is outside OFFICE_WORKSPACE_ROOT")
                rel_dst = os.path.relpath(dst, workspace_root)
                raw_dst = os.path.join(workspace_root, _normalize_organized_path(rel_dst))
            else:
                raw_dst = os.path.join(workspace_root, _normalize_organized_path(dst))
            dst_normalized, err = _validate_workspace_path(raw_dst)
            if err:
                return ToolResult(output="", error=f"organize_move_file: destination {err}")

        # In workspace mode, allow organizing into subdirectories under organized-output/files/
        # The validation is just to ensure we don't escape OFFICE_WORKSPACE_ROOT

        # Validate source path if provided
        if src:
            src_normalized, err = _validate_path(src)
            if err:
                return ToolResult(output="", error=f"organize_move_file: source {err}")
        else:
            src_normalized = ""

        if action == "copy_file" and src_normalized:
            try:
                metadata_root = _metadata_root_for_path(src_normalized)
                source_metadata = _build_file_metadata(metadata_root, src_normalized)
            except Exception:
                source_metadata = {}
            expected_entity = _safe_path_segment(str(source_metadata.get("primary_entity") or ""))
            expected_date = str(source_metadata.get("inferred_date_bucket") or "")
            expected_filename = str(source_metadata.get("relative_path") or "").replace(os.sep, "-")
            confidence = str(source_metadata.get("primary_entity_confidence") or "")
            if confidence == "high" and expected_entity and expected_date and expected_filename:
                dst_parts = Path(dst_normalized).parts
                tail = list(dst_parts[-3:]) if len(dst_parts) >= 3 else list(dst_parts)
                expected_tail = [expected_entity, expected_date, expected_filename]
                if tail != expected_tail:
                    return ToolResult(
                        output="",
                        error=(
                            "organize_move_file: destination does not match high-confidence source metadata. "
                            f"Expected tail {'/'.join(expected_tail)!r}, got {'/'.join(tail)!r}"
                        ),
                    )

        plan_path = os.path.join(workspace_root, "operations-plan.json") if workspace_root else ""
        if action == "copy_file" and src_normalized and plan_path and os.path.exists(plan_path):
            try:
                with open(plan_path, encoding="utf-8") as existing_plan:
                    for raw_line in existing_plan:
                        raw_line = raw_line.strip()
                        if not raw_line:
                            continue
                        record = json.loads(raw_line)
                        if (
                            record.get("action") == "copy_file"
                            and record.get("status") == "succeeded"
                            and record.get("src") == src_normalized
                        ):
                            return ToolResult(
                                output="",
                                error=(
                                    "organize_move_file: source file already copied successfully in this task. "
                                    f"Duplicate copy blocked for {src_normalized}"
                                ),
                            )
            except Exception:
                pass

        # Execute action
        try:
            if action == "mkdir":
                os.makedirs(dst_normalized, exist_ok=True)
                if plan_path:
                    with open(plan_path, "a", encoding="utf-8") as f:
                        f.write(json.dumps({
                            "action": action,
                            "src": src_normalized,
                            "dst": dst_normalized,
                            "content_length": len(content) if content else 0,
                            "status": "succeeded",
                        }) + "\n")
                return ToolResult(output=json.dumps({"path": dst_normalized, "action": "mkdir"}))
            elif action == "copy_file":
                if not src_normalized:
                    return ToolResult(output="", error="organize_move_file: copy_file requires src")
                os.makedirs(os.path.dirname(dst_normalized), exist_ok=True)
                shutil.copy2(src_normalized, dst_normalized)
                if plan_path:
                    with open(plan_path, "a", encoding="utf-8") as f:
                        f.write(json.dumps({
                            "action": action,
                            "src": src_normalized,
                            "dst": dst_normalized,
                            "content_length": len(content) if content else 0,
                            "status": "succeeded",
                        }) + "\n")
                return ToolResult(output=json.dumps({"from": src_normalized, "to": dst_normalized, "action": "copy_file"}))
            elif action == "write_text":
                if not content:
                    return ToolResult(output="", error="organize_move_file: write_text requires content")
                os.makedirs(os.path.dirname(dst_normalized), exist_ok=True)
                with open(dst_normalized, "w", encoding="utf-8") as f:
                    f.write(content)
                if plan_path:
                    with open(plan_path, "a", encoding="utf-8") as f:
                        f.write(json.dumps({
                            "action": action,
                            "src": src_normalized,
                            "dst": dst_normalized,
                            "content_length": len(content) if content else 0,
                            "status": "succeeded",
                        }) + "\n")
                return ToolResult(output=json.dumps({"path": dst_normalized, "action": "write_text", "bytes": len(content)}))
        except Exception as exc:
            if plan_path:
                try:
                    with open(plan_path, "a", encoding="utf-8") as f:
                        f.write(json.dumps({
                            "action": action,
                            "src": src_normalized,
                            "dst": dst_normalized,
                            "content_length": len(content) if content else 0,
                            "status": "failed",
                            "error": str(exc),
                        }) + "\n")
                except Exception:
                    pass
            return ToolResult(output="", error=f"organize_move_file: {exc}")


# ---------------------------------------------------------------------------
# Tool registration
# ---------------------------------------------------------------------------

_OFFICE_TOOLS = [
    ReadPdfTool(),
    ReadDocxTool(),
    ReadPptxTool(),
    ReadTxtTool(),
    ReadCsvTool(),
    ReadXlsxTool(),
    ReadXlsTool(),
    ListDirectoryTool(),
    WriteWorkspaceTool(),
    WriteFileTool(),
    OrganizeFolderTool(),
    OrganizeMoveFileTool(),
]


def register_office_tools() -> None:
    """Register office tools into the global ToolRegistry (idempotent)."""
    from framework.tools.registry import get_registry
    registry = get_registry()
    for tool in _OFFICE_TOOLS:
        registry.register(tool)


# ---------------------------------------------------------------------------
# Delete Output File Tool
# ---------------------------------------------------------------------------

class DeleteOutputFileTool(BaseTool):
    """Delete a file under the resolved task output root.

    Hard rules (enforced in this order, fail-closed):
    1. ``OFFICE_RESOLVED_TARGET_DIR`` (or ``OFFICE_WORKSPACE_ROOT`` as a
       fallback for workspace mode) must be set to a real directory; the
       tool refuses otherwise.
    2. The candidate path resolves with ``realpath`` (following symlinks).
    3. The candidate path must NOT resolve to a file under any validated
       source root (``OFFICE_SOURCE_ROOT``) — checked first so a symlink
       escape into a source folder is reported as a source violation.
    4. The candidate path must lie inside the resolved target directory.
    5. Only regular files are removed; non-regular targets and missing
       files are refused with a descriptive error.
    """
    name = "delete_output_file"
    description = "Delete a stale output file under the resolved task output root. Refuses source inputs and out-of-root paths."
    parameters_schema = {
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "Path of the file to delete, relative to the resolved target directory or absolute under it.",
            },
        },
        "required": ["filename"],
    }

    def execute_sync(self, filename: str = "") -> ToolResult:
        target_dir = os.environ.get("OFFICE_RESOLVED_TARGET_DIR", "").strip()
        if not target_dir:
            target_dir = os.environ.get("OFFICE_WORKSPACE_ROOT", "").strip()
        if not target_dir:
            return ToolResult(output="", error="delete_output_file: OFFICE_RESOLVED_TARGET_DIR (or OFFICE_WORKSPACE_ROOT) is not set")
        try:
            real_target = os.path.realpath(os.path.abspath(target_dir))
        except Exception as exc:
            return ToolResult(output="", error=f"delete_output_file: target resolution failed: {exc}")
        if not os.path.isdir(real_target):
            return ToolResult(output="", error=f"delete_output_file: target {real_target!r} is not a directory")
        if not filename:
            return ToolResult(output="", error="delete_output_file: filename is required")
        # candidate resolution — follow symlinks via realpath so escape
        # attempts (including symlink chains) are caught by the next two
        # checks rather than slipping through with the link's path.
        try:
            if os.path.isabs(filename):
                candidate = os.path.realpath(os.path.abspath(filename))
            else:
                candidate = os.path.realpath(os.path.join(real_target, filename))
        except Exception as exc:
            return ToolResult(output="", error=f"delete_output_file: path resolution failed: {exc}")
        # source-input protection — must precede the target-containment
        # check so a symlink into a source folder reports as a source
        # violation rather than a generic "outside target" error.
        source_root = os.environ.get("OFFICE_SOURCE_ROOT", "").strip()
        if source_root:
            real_source = os.path.realpath(os.path.abspath(source_root))
            if candidate == real_source or candidate.startswith(real_source.rstrip(os.sep) + os.sep):
                return ToolResult(output="", error=f"delete_output_file: refusing to delete source input {filename!r}")
        if not (candidate == real_target or candidate.startswith(real_target.rstrip(os.sep) + os.sep)):
            return ToolResult(output="", error=f"delete_output_file: path {filename!r} is outside the resolved target directory")
        if not os.path.exists(candidate):
            return ToolResult(output="", error=f"delete_output_file: file does not exist: {filename!r}")
        if not os.path.isfile(candidate):
            return ToolResult(output="", error=f"delete_output_file: refusing to delete non-regular file: {filename!r}")
        try:
            os.remove(candidate)
        except OSError as exc:
            return ToolResult(output="", error=f"delete_output_file: remove failed: {exc}")
        return ToolResult(output=json.dumps({"deleted": candidate}))
