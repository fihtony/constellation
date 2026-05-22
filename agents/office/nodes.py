"""Office agent workflow nodes.

receive_task     — Parse task message: capability, source paths, output mode
analyze_request — Validate paths, check permissions, load skill prompts
execute_office_work — ReAct core: runtime.run_agentic() with office tools
report_result   — Write pr-evidence.json, return result
"""

from __future__ import annotations

import json
import logging
import os
import re
import csv
import threading
from typing import Any

from agents.office.office_tools import (
    _check_directory_limits,
    collect_organize_file_inventory,
)

logger = logging.getLogger(__name__)

AGENT_ID = "office"
SUMMARY_EXTENSIONS = {
    ".pdf", ".docx", ".docm", ".dotx", ".dotm",
    ".txt", ".md", ".markdown", ".html", ".htm", ".xml",
    ".json", ".jsonl", ".yaml", ".yml", ".log", ".ini", ".cfg", ".toml", ".rtf",
    ".pptx", ".csv", ".tsv", ".xlsx", ".xls",
}


def _contains_cjk(text: str) -> bool:
    return any(
        "\u4e00" <= ch <= "\u9fff" or
        "\u3400" <= ch <= "\u4dbf" or
        "\u3040" <= ch <= "\u30ff" or
        "\uac00" <= ch <= "\ud7af"
        for ch in text
    )


def _english_summary_for_report(capability: str, success: bool, expected_outputs: list[str]) -> str:
    action_map = {
        "summarize": "document summarization",
        "analyze": "data analysis",
        "organize": "folder organization",
    }
    action = action_map.get(capability, "office work")
    if success:
        if expected_outputs:
            outputs = ", ".join(os.path.basename(path) for path in expected_outputs[:5])
            return f"Office {action} completed successfully. Primary outputs: {outputs}."
        return f"Office {action} completed successfully."
    return f"Office {action} did not complete successfully. Review warnings and task artifacts for details."


def _english_agentic_output(raw_output: str, capability: str, expected_outputs: list[str]) -> str:
    if raw_output and not _contains_cjk(raw_output):
        return raw_output
    lines = [
        "The original agentic response was not persisted verbatim because it was not fully in English.",
        _english_summary_for_report(capability, True, expected_outputs),
    ]
    if expected_outputs:
        lines.append("Generated outputs:")
        lines.extend(f"- {path}" for path in expected_outputs)
    return "\n".join(lines) + "\n"


# ---------------------------------------------------------------------------
# Path helpers
# ---------------------------------------------------------------------------

def _get_workspace_root(state: dict) -> str:
    """Get the workspace root for this task.

    Workspace path: {ARTIFACT_ROOT}/{compass_task_id}/office/
    All office tasks under the same compass task share the same workspace.
    """
    artifact_root = os.environ.get(
        "ARTIFACT_ROOT",
        os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "artifacts")
    )
    compass_id = state.get("_compass_task_id", "default")
    return os.path.join(artifact_root, compass_id, "office")


def _validate_source_path(path: str) -> tuple[str, str]:
    """Validate path is within OFFICE_SOURCE_ROOT."""
    source_root = os.environ.get("OFFICE_SOURCE_ROOT", "/")
    real_path = os.path.realpath(os.path.abspath(path))
    real_root = os.path.realpath(os.path.abspath(source_root))
    prefix = real_root.rstrip(os.sep) + os.sep
    if real_path != real_root and not real_path.startswith(prefix):
        return "", f"Path {path!r} is outside OFFICE_SOURCE_ROOT ({source_root})"
    return real_path, ""


# ---------------------------------------------------------------------------
# Capability parsing
# ---------------------------------------------------------------------------

def _parse_capability(text: str) -> str:
    """Infer capability from task text."""
    text_lower = text.lower()
    if any(kw in text_lower for kw in ("summarize", "summary", "summaries")):
        return "summarize"
    if "analyze" in text_lower and any(
        kw in text_lower for kw in ("csv", "data", "xlsx", "xls", "spreadsheet", "table", "report")
    ):
        return "analyze"
    if "organize" in text_lower or "folder" in text_lower:
        return "organize"
    return "summarize"


def _extract_paths(text: str) -> list[str]:
    """Extract file/folder paths from task text."""
    # Match absolute Unix paths
    paths = re.findall(r'(?:^|\s)((?:/[a-zA-Z0-9_.~-]+)+)', text)
    # Also match relative paths
    rel_paths = re.findall(r'(?:^|\s)([a-zA-Z0-9_./-]*(?:tests?/data[/\w.-]+)[a-zA-Z0-9_./-]*)', text)
    all_paths = paths + rel_paths
    def _sanitize(candidate: str) -> str:
        return candidate.strip().strip("\"'`").lstrip("([{").rstrip(".,;:!?)]}\"'`")
    # Deduplicate while preserving order
    normalized = []
    for p in all_paths:
        sp = _sanitize(p)
        if len(sp) > 3:
            normalized.append(sp)
    return list(dict.fromkeys(normalized))


def _expand_summarize_sources(paths: list[str]) -> list[str]:
    """Expand directory inputs into supported document files for summarize tasks."""
    expanded: list[str] = []
    for path in paths:
        if os.path.isdir(path):
            for root, _, files in os.walk(path):
                for name in sorted(files):
                    if os.path.splitext(name)[1].lower() in SUMMARY_EXTENSIONS:
                        expanded.append(os.path.join(root, name))
        else:
            expanded.append(path)
    return list(dict.fromkeys(expanded))


# ---------------------------------------------------------------------------
# Node: receive_task
# ---------------------------------------------------------------------------

def receive_task(state: dict) -> dict:
    """Parse the incoming task message to extract capability, paths, and output mode."""
    user_text = state.get("user_request", "")

    # Parse output mode
    output_mode = "workspace"
    if "inplace" in user_text.lower():
        output_mode = "inplace"

    # Parse capability
    capability = _parse_capability(user_text)

    # Parse source paths
    source_paths = _extract_paths(user_text)

    logger.info(f"receive_task: capability={capability} output_mode={output_mode} paths={source_paths}")

    return {
        "capability": capability,
        "output_mode": output_mode,
        "source_paths": source_paths,
    }


# ---------------------------------------------------------------------------
# Node: analyze_request
# ---------------------------------------------------------------------------

def analyze_request(state: dict) -> dict:
    """Validate source paths and check permissions. Returns updated state or error."""
    source_paths = state.get("source_paths", [])
    output_mode = state.get("output_mode", "workspace")
    capability = state.get("capability", "summarize")

    if not source_paths:
        return {"error": "No source paths found in task. Please provide file or folder paths."}

    validated_paths = []
    for p in source_paths:
        normalized, err = _validate_source_path(p)
        if normalized:
            validated_paths.append(normalized)
        else:
            logger.warning(f"Skipping invalid path: {p} — {err}")

    if capability == "summarize":
        validated_paths = _expand_summarize_sources(validated_paths)

    if not validated_paths and capability not in ("summarize", "organize"):
        return {"error": "No valid paths found under OFFICE_SOURCE_ROOT."}

    # Directory resource pre-check for organize capability
    if capability == "organize" and validated_paths:
        # Check first path (assuming single directory for organize)
        first_path = validated_paths[0]
        if os.path.isdir(first_path):
            limit_error = _check_directory_limits(first_path)
            if limit_error:
                return limit_error

    workspace_root = _get_workspace_root(state)
    os.makedirs(workspace_root, exist_ok=True)
    artifacts_dir = os.path.join(workspace_root, "artifacts")
    os.makedirs(artifacts_dir, exist_ok=True)

    # Check inplace permission
    if output_mode == "inplace":
        allow_inplace = os.environ.get("OFFICE_ALLOW_INPLACE_WRITES", "false").lower()
        if allow_inplace not in ("true", "1", "yes"):
            logger.warning("inplace mode requested but OFFICE_ALLOW_INPLACE_WRITES not set — falling back to workspace")
            state["output_mode"] = "workspace"
            output_mode = "workspace"

    os.environ["OFFICE_OUTPUT_MODE"] = output_mode

    logger.info(f"analyze_request: validated_paths={validated_paths} artifacts_dir={artifacts_dir}")

    return {
        "validated_paths": validated_paths,
        "workspace_root": workspace_root,
        "artifacts_dir": artifacts_dir,
    }


# ---------------------------------------------------------------------------
# System prompt loader
# ---------------------------------------------------------------------------

def _load_system_prompt() -> str:
    """Load the office agent system prompt."""
    prompt_path = os.path.join(os.path.dirname(__file__), "prompts", "system.md")
    if os.path.exists(prompt_path):
        with open(prompt_path, encoding="utf-8") as fh:
            return fh.read()
    return (
        "You are an expert office task agent. "
        "Use the read_pdf, read_docx, read_txt, read_csv tools to read files. "
        "Use write_workspace to save results to the workspace artifacts folder. "
        "Be concise and thorough."
    )


def _build_skill_context(state: dict) -> str:
    """Build prompt context from configured office skills."""
    skills_registry = state.get("_skills_registry")
    required_skills = state.get("required_skills", [])
    if not skills_registry or not required_skills:
        return ""
    try:
        return skills_registry.build_prompt_context(required_skills)
    except Exception:
        return ""


def _capability_tool_names(capability: str, output_mode: str) -> list[str]:
    """Return the minimal MCP tool surface for the current office task."""
    if capability == "analyze":
        tools = [
            "read_csv",
            "read_txt",
            "read_xlsx",
            "read_xls",
            "read_pdf",
            "read_docx",
            "list_directory",
        ]
        tools.append("write_file" if output_mode == "inplace" else "write_workspace")
        return tools

    if capability == "summarize":
        tools = [
            "read_pdf",
            "read_docx",
            "read_txt",
            "read_pptx",
            "read_csv",
            "read_xlsx",
            "read_xls",
            "list_directory",
        ]
        tools.append("write_file" if output_mode == "inplace" else "write_workspace")
        return tools

    if capability == "organize":
        tools = [
            "list_directory",
            "organize_folder",
            "organize_move_file",
            "read_txt",
            "read_pdf",
            "read_docx",
            "read_pptx",
            "read_csv",
            "read_xlsx",
            "read_xls",
        ]
        tools.append("write_file" if output_mode == "inplace" else "write_workspace")
        return tools

    return []


def _claude_allowed_tool_names(tool_names: list[str]) -> list[str]:
    """Map Constellation tool ids to Claude Code MCP tool ids."""
    return [f"mcp__constellation_tools__{tool_name}" for tool_name in tool_names]


def _expected_output_paths(
    capability: str,
    validated_paths: list[str],
    output_mode: str,
    artifacts_dir: str,
) -> list[str]:
    """Return the required delivery files for the current office task."""
    expected: list[str] = []
    if capability == "analyze":
        for path in validated_paths:
            if os.path.isfile(path):
                expected.append(_target_output_path(output_mode, path, artifacts_dir, ".analysis.md"))
    elif capability == "summarize":
        for path in validated_paths:
            if os.path.isfile(path):
                expected.append(_target_output_path(output_mode, path, artifacts_dir, ".summary.md"))
        if len([path for path in validated_paths if os.path.isfile(path)]) > 1:
            base_path = validated_paths[0]
            expected.append(_target_output_file(output_mode, base_path, artifacts_dir, "combined-summary.md"))
    elif capability == "organize" and validated_paths:
        expected.append(_target_output_file(output_mode, validated_paths[0], artifacts_dir, "organization-plan.md"))
    return expected


def _verify_delivery_paths(expected_paths: list[str], output_mode: str, artifacts_dir: str) -> tuple[bool, list[str]]:
    """Check that all required outputs exist and stay inside the authorized area."""
    errors: list[str] = []
    workspace_root = os.path.realpath(os.path.abspath(artifacts_dir)) if artifacts_dir else ""
    for path in expected_paths:
        real_path = os.path.realpath(os.path.abspath(path))
        if not os.path.exists(real_path):
            errors.append(f"Missing expected output: {path}")
            continue
        if output_mode == "workspace" and workspace_root:
            prefix = workspace_root.rstrip(os.sep) + os.sep
            if real_path != workspace_root and not real_path.startswith(prefix):
                errors.append(f"Output escaped workspace root: {real_path}")
    return (not errors), errors


def _verify_organize_materialization(output_mode: str, artifacts_dir: str, source_paths: list[str]) -> list[str]:
    """Ensure organize tasks created a real organized-output tree, not only a plan."""
    if output_mode == "workspace":
        root = os.path.join(artifacts_dir, "organized-output", "files")
    else:
        source_root = source_paths[0] if source_paths else ""
        root = os.path.join(source_root, "organized-output", "files")
    if not os.path.isdir(root):
        return [f"Missing organized output directory: {root}"]
    copied_files: list[str] = []
    for walk_root, _, files in os.walk(root):
        for name in files:
            copied_files.append(os.path.join(walk_root, name))
    if not copied_files:
        return [f"No organized files were materialized under: {root}"]
    if not source_paths:
        return []

    operations_path = os.path.join(artifacts_dir, "operations-plan.json")
    if not os.path.exists(operations_path):
        return [f"Missing operations log: {operations_path}"]

    source_root = source_paths[0]
    inventory, _, _ = collect_organize_file_inventory(source_root)
    expected_sources = {
        os.path.realpath(os.path.join(source_root, str(item["relative_path"])))
        for item in inventory
    }

    copy_actions: list[dict[str, str]] = []
    try:
        with open(operations_path, encoding="utf-8") as fh:
            for raw_line in fh:
                line = raw_line.strip()
                if not line:
                    continue
                record = json.loads(line)
                if record.get("action") == "copy_file" and record.get("src") and record.get("dst"):
                    if record.get("status") == "failed":
                        continue
                    copy_actions.append({
                        "src": os.path.realpath(os.path.abspath(record["src"])),
                        "dst": os.path.realpath(os.path.abspath(record["dst"])),
                    })
    except Exception as exc:
        return [f"Failed to read operations log {operations_path}: {exc}"]

    if not copy_actions:
        return [f"No copy_file operations recorded in {operations_path}"]

    copy_counts: dict[str, int] = {}
    for record in copy_actions:
        src = record["src"]
        copy_counts[src] = copy_counts.get(src, 0) + 1

    errors: list[str] = []
    duplicated = sorted(src for src, count in copy_counts.items() if count > 1)
    if duplicated:
        errors.append(
            "Source files copied more than once: "
            + ", ".join(os.path.relpath(src, source_root) for src in duplicated[:10])
        )

    missing = sorted(src for src in expected_sources if copy_counts.get(src, 0) == 0)
    if missing:
        errors.append(
            "Source files were not copied: "
            + ", ".join(os.path.relpath(src, source_root) for src in missing[:10])
        )

    def _path_key(value: str) -> str:
        return re.sub(r"[^a-z0-9]+", "", value.lower())

    inventory_by_src = {
        os.path.realpath(os.path.join(source_root, str(item["relative_path"]))): item
        for item in inventory
    }
    for record in copy_actions:
        if not os.path.exists(record["dst"]):
            continue
        item = inventory_by_src.get(record["src"])
        if not item:
            continue
        dst_relative = os.path.relpath(record["dst"], root)
        primary_entity = str(item.get("primary_entity") or "").strip()
        primary_entity_confidence = str(item.get("primary_entity_confidence") or "")
        if primary_entity and primary_entity_confidence == "high" and _path_key(primary_entity) not in _path_key(dst_relative):
            errors.append(
                f"Destination does not reflect inferred entity for {item['relative_path']}: "
                f"{primary_entity} -> {dst_relative}"
            )
        date_bucket = str(item.get("inferred_date_bucket") or "").strip()
        if date_bucket and date_bucket not in dst_relative:
            errors.append(
                f"Destination does not reflect inferred date bucket for {item['relative_path']}: "
                f"{date_bucket} -> {dst_relative}"
            )

    return errors


def _effective_agentic_budget(capability: str, validated_paths: list[str]) -> tuple[int, int]:
    """Scale agentic budget with workload size instead of using a single flat timeout."""
    max_turns = int(os.environ.get("OFFICE_AGENTIC_MAX_TURNS", "30"))
    timeout_seconds = int(os.environ.get("OFFICE_AGENTIC_TIMEOUT_SECONDS", "1800"))
    if capability == "summarize":
        doc_count = len([path for path in validated_paths if os.path.isfile(path)])
        if doc_count > 1:
            timeout_seconds = max(timeout_seconds, min(900, 120 + doc_count * 45))
            max_turns = max(max_turns, min(40, 8 + doc_count * 2))
    elif capability == "organize":
        folder_count = len([path for path in validated_paths if os.path.isdir(path)])
        base_timeout = 900
        if folder_count == 1:
            subdirs = 0
            for p in validated_paths:
                if os.path.isdir(p):
                    subdirs = len([d for d in os.listdir(p) if os.path.isdir(os.path.join(p, d)) and not d.startswith('.')])
            timeout_seconds = max(base_timeout, min(1200, base_timeout + subdirs * 25))
            max_turns = max(max_turns, min(60, 30 + subdirs * 2))
        else:
            timeout_seconds = max(base_timeout, min(1200, base_timeout + folder_count * 30))
    return max_turns, timeout_seconds


# ---------------------------------------------------------------------------
# Node: execute_office_work
# ---------------------------------------------------------------------------

def execute_office_work(state: dict) -> dict:
    """ReAct core: call runtime.run_agentic() with office tools to do the actual work."""
    runtime = state.get("_runtime")
    if not runtime:
        return {"error": "No runtime configured"}

    capability = state.get("capability", "summarize")
    validated_paths = state.get("validated_paths", [])
    artifacts_dir = state.get("artifacts_dir", "")
    output_mode = state.get("output_mode", "workspace")
    source_root = os.environ.get("OFFICE_SOURCE_ROOT", "")

    if capability == "organize":
        prompt = _build_organize_prompt(validated_paths, output_mode, source_root)
    elif capability == "summarize":
        prompt = _build_summarize_prompt(validated_paths, output_mode, source_root)
    elif capability == "analyze":
        prompt = _build_analyze_prompt(validated_paths, output_mode, source_root)
    else:
        return {"error": f"Capability {capability} not implemented."}

    # Set workspace root env for write_workspace tool
    if artifacts_dir:
        os.environ["OFFICE_WORKSPACE_ROOT"] = artifacts_dir

    logger.info(f"execute_office_work: running {capability} with {len(validated_paths)} paths")

    tool_names = _capability_tool_names(capability, output_mode)
    if not tool_names:
        return {"error": f"No tools configured for capability {capability!r}"}

    max_turns, timeout_seconds = _effective_agentic_budget(capability, validated_paths)

    skill_context = _build_skill_context(state)
    system_prompt = _load_system_prompt()
    if skill_context:
        system_prompt = (
            f"{system_prompt}\n\n"
            "## Loaded Skill Context\n"
            "The following methodology skills are mandatory for this run.\n\n"
            f"{skill_context}"
        )

    def _run_agentic_call():
        return runtime.run_agentic(
            prompt,
            system_prompt=system_prompt,
            cwd=state.get("workspace_root") or (
                validated_paths[0] if validated_paths and os.path.isdir(validated_paths[0]) else (
                    os.path.dirname(validated_paths[0]) if validated_paths else None
                )
            ),
            tools=tool_names,
            allowed_tools=_claude_allowed_tool_names(tool_names),
            max_turns=max_turns,
            timeout=timeout_seconds,
            plugin_manager=state.get("_plugin_manager"),
        )

    result_holder: dict[str, Any] = {}
    error_holder: dict[str, str] = {}

    def _worker():
        try:
            result_holder["result"] = _run_agentic_call()
        except Exception as exc:
            error_holder["error"] = str(exc)

    worker = threading.Thread(target=_worker, daemon=True)
    worker.start()
    worker.join(timeout=timeout_seconds + 15)

    if worker.is_alive():
        from framework.runtime.adapter import AgenticResult
        result = AgenticResult(
            success=False,
            summary=f"agentic runtime watchdog timeout after {timeout_seconds + 15}s",
            backend_used="watchdog-timeout",
        )
    elif error_holder.get("error"):
        from framework.runtime.adapter import AgenticResult
        result = AgenticResult(
            success=False,
            summary=f"agentic runtime error: {error_holder.get('error')}",
            backend_used="watchdog-error",
        )
    else:
        result = result_holder.get("result")

    if result.success:
        expected_outputs = _expected_output_paths(capability, validated_paths, output_mode, artifacts_dir)
        delivery_ok, delivery_errors = _verify_delivery_paths(expected_outputs, output_mode, artifacts_dir)
        if capability == "organize":
            delivery_errors.extend(_verify_organize_materialization(output_mode, artifacts_dir, validated_paths))
            delivery_ok = not delivery_errors
        if not delivery_ok:
            logger.error("execute_office_work: delivery verification failed: %s", "; ".join(delivery_errors))
            return {
                "summary": (
                    "Agentic execution completed, but delivery verification failed.\n"
                    + "\n".join(f"- {err}" for err in delivery_errors)
                ),
                "success": False,
                "capability": capability,
                "status": "failed",
                "raw_output": getattr(result, "raw_output", ""),
                "warnings": delivery_errors,
                "error": "delivery verification failed",
            }
        logger.info(f"execute_office_work: success")
        return {
            "summary": result.summary,
            "success": True,
            "capability": capability,
            "status": "completed",
            "raw_output": getattr(result, "raw_output", ""),
            "expected_outputs": expected_outputs,
        }

    logger.error(f"execute_office_work: failed — {result.summary}")
    return {
        "summary": result.summary,
        "success": False,
        "capability": capability,
        "status": "failed",
        "raw_output": getattr(result, "raw_output", ""),
        "warnings": [f"Agentic runtime failed: {result.summary}."],
        "error": result.summary,
    }


def _target_output_file(output_mode: str, source_path: str, artifacts_dir: str, filename: str) -> str:
    if output_mode == "inplace":
        base_dir = source_path if os.path.isdir(source_path) else os.path.dirname(source_path)
        return os.path.join(base_dir, os.path.basename(filename))
    return os.path.join(artifacts_dir, os.path.basename(filename))


def _target_output_path(output_mode: str, source_path: str, artifacts_dir: str, suffix: str) -> str:
    basename = os.path.basename(source_path.rstrip("/"))
    return _target_output_file(output_mode, source_path, artifacts_dir, f"{basename}{suffix}")


def _maybe_float(value: str) -> float | None:
    try:
        cleaned = str(value).strip().replace(",", "")
        if not cleaned:
            return None
        return float(cleaned)
    except (TypeError, ValueError):
        return None


def _detect_numeric_columns(rows: list[dict[str, str]], headers: list[str]) -> set[str]:
    numeric_cols: set[str] = set()
    for h in headers:
        non_empty = 0
        numeric = 0
        for row in rows:
            val = row.get(h, "")
            if str(val).strip() == "":
                continue
            non_empty += 1
            if _maybe_float(val) is not None:
                numeric += 1
        if non_empty > 0 and numeric / non_empty >= 0.8:
            numeric_cols.add(h)
    return numeric_cols


def _load_tabular_rows(path: str) -> tuple[list[str], list[dict[str, str]], str]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".csv":
        with open(path, encoding="utf-8", errors="replace") as fh:
            reader = csv.DictReader(fh)
            rows = list(reader)
            return list(reader.fieldnames or []), rows, "csv"
    if ext in {".xlsx", ".xlsm"}:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        ws = wb[wb.sheetnames[0]]
        raw_rows = [list(r) for r in ws.iter_rows(values_only=True)]
        if not raw_rows:
            return [], [], "xlsx"
        headers = [str(h).strip() if h is not None else f"column_{idx+1}" for idx, h in enumerate(raw_rows[0])]
        rows: list[dict[str, str]] = []
        for raw in raw_rows[1:]:
            row = {headers[i]: ("" if i >= len(raw) or raw[i] is None else str(raw[i])) for i in range(len(headers))}
            rows.append(row)
        return headers, rows, "xlsx"
    if ext == ".xls":
        import xlrd
        wb = xlrd.open_workbook(path)
        ws = wb.sheet_by_index(0)
        if ws.nrows == 0:
            return [], [], "xls"
        headers = [str(ws.cell_value(0, i)).strip() or f"column_{i+1}" for i in range(ws.ncols)]
        rows: list[dict[str, str]] = []
        for r in range(1, ws.nrows):
            row = {headers[c]: str(ws.cell_value(r, c)) for c in range(ws.ncols)}
            rows.append(row)
        return headers, rows, "xls"
    if ext in {".txt", ".tsv"}:
        with open(path, encoding="utf-8", errors="replace") as fh:
            sample = fh.read(4096)
            fh.seek(0)
            dialect = csv.excel
            try:
                dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|")
            except Exception:
                pass
            reader = csv.DictReader(fh, dialect=dialect)
            rows = list(reader)
            return list(reader.fieldnames or []), rows, "txt"
    return [], [], "unsupported"


def _extract_document_text(path: str, max_chars: int = 10000) -> str:
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext in {".txt", ".md", ".markdown", ".json", ".jsonl", ".yaml", ".yml", ".log", ".ini", ".cfg", ".toml", ".rtf"}:
            with open(path, encoding="utf-8", errors="replace") as fh:
                raw = fh.read(max_chars * 2)
            from agents.office.office_tools import _extract_markup_text
            text, _ = _extract_markup_text(raw, ext)
            return text[:max_chars]
        if ext in {".csv", ".tsv"}:
            headers, rows, _ = _load_tabular_rows(path)
            preview = {
                "headers": headers,
                "sample_rows": rows[:10],
            }
            return json.dumps(preview, ensure_ascii=False)[:max_chars]
        if ext in {".html", ".htm", ".xml"}:
            from agents.office.office_tools import _extract_markup_text
            with open(path, encoding="utf-8", errors="replace") as fh:
                raw = fh.read(max_chars * 2)
            text, _ = _extract_markup_text(raw, ext)
            return text[:max_chars]
        if ext == ".pdf":
            import pdfplumber
            parts = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    txt = page.extract_text() or ""
                    if txt.strip():
                        parts.append(txt.strip())
                    if sum(len(p) for p in parts) >= max_chars:
                        break
            return "\n\n".join(parts)[:max_chars]
        if ext in {".docx", ".docm", ".dotx", ".dotm"}:
            from agents.office.office_tools import _extract_docx_like_text
            lines, _ = _extract_docx_like_text(path)
            return "\n".join(lines)[:max_chars]
        if ext == ".xlsx":
            headers, rows, _ = _load_tabular_rows(path)
            preview = {
                "headers": headers,
                "sample_rows": rows[:10],
            }
            return json.dumps(preview, ensure_ascii=False)[:max_chars]
        if ext == ".xls":
            headers, rows, _ = _load_tabular_rows(path)
            preview = {
                "headers": headers,
                "sample_rows": rows[:10],
            }
            return json.dumps(preview, ensure_ascii=False)[:max_chars]
        if ext == ".pptx":
            import pptx
            prs = pptx.Presentation(path)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        txt = shape.text.strip()
                        if txt:
                            lines.append(txt)
                if sum(len(x) for x in lines) >= max_chars:
                    break
            return "\n".join(lines)[:max_chars]
    except Exception:
        return ""
    return ""


def _make_bullet_points(text: str, max_points: int = 5) -> list[str]:
    lines = [ln.strip(" -\t") for ln in text.splitlines() if ln.strip()]
    bullets: list[str] = []
    seen: set[str] = set()
    for ln in lines:
        key = ln.lower()
        if key in seen:
            continue
        seen.add(key)
        bullets.append(ln[:180])
        if len(bullets) >= max_points:
            break
    return bullets


def _write_text(path: str, content: str) -> None:
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "w", encoding="utf-8") as fh:
        fh.write(content)


def _fallback_analyze(paths: list[str], output_mode: str, artifacts_dir: str) -> dict:
    table_paths = [
        p for p in paths
        if os.path.isfile(p) and os.path.splitext(p)[1].lower() in {".csv", ".xlsx", ".xlsm", ".xls", ".txt", ".tsv"}
    ]
    if not table_paths:
        return {"success": False, "summary": "No readable tabular path provided for fallback analysis."}
    outputs: list[str] = []
    for table_path in table_paths:
        headers, rows, source_type = _load_tabular_rows(table_path)
        numeric_cols = _detect_numeric_columns(rows, headers)
        numeric_stats: dict[str, dict[str, float]] = {}
        for col in sorted(numeric_cols):
            vals = [_maybe_float(r.get(col, "")) for r in rows]
            nums = [v for v in vals if v is not None]
            if not nums:
                continue
            numeric_stats[col] = {
                "count": float(len(nums)),
                "min": min(nums),
                "max": max(nums),
                "avg": sum(nums) / len(nums),
            }
        categorical_cols = [h for h in headers if h not in numeric_cols]
        top_group_lines: list[str] = []
        if categorical_cols and numeric_cols:
            for cat_col in categorical_cols[:3]:
                for num_col in sorted(numeric_cols)[:3]:
                    grouped: dict[str, float] = {}
                    for row in rows:
                        key = (row.get(cat_col, "") or "Unknown").strip() or "Unknown"
                        val = _maybe_float(row.get(num_col, ""))
                        if val is None:
                            continue
                        grouped[key] = grouped.get(key, 0.0) + val
                    if not grouped:
                        continue
                    top_items = sorted(grouped.items(), key=lambda kv: kv[1], reverse=True)[:3]
                    preview = ", ".join(f"{k}={v:.2f}" for k, v in top_items)
                    top_group_lines.append(f"- `{cat_col}` grouped by `{num_col}` (top 3): {preview}")
        lines = [
            f"# Data Analysis: {os.path.basename(table_path)}",
            "",
            "## File Overview",
            f"- Source type: {source_type}",
            f"- Rows: {len(rows)}",
            f"- Columns: {', '.join(headers) if headers else 'N/A'}",
            "",
            "## Numeric Summary",
        ]
        if numeric_stats:
            for col, st in numeric_stats.items():
                lines.append(
                    f"- `{col}`: count={int(st['count'])}, min={st['min']:.2f}, max={st['max']:.2f}, avg={st['avg']:.2f}"
                )
        else:
            lines.append("- No stable numeric columns detected.")
        lines.extend(["", "## Key Insights"])
        if top_group_lines:
            lines.extend(top_group_lines)
        else:
            lines.append("- No reliable categorical+numeric aggregation insight detected.")
        lines.extend([
            "- Insights were inferred from detected schema, without hardcoded column assumptions.",
            "- Report generated by deterministic fallback after agentic runtime failure.",
            "",
        ])
        out_path = _target_output_path(output_mode, table_path, artifacts_dir, ".analysis.md")
        _write_text(out_path, "\n".join(lines))
        outputs.append(out_path)
    return {"success": True, "summary": f"Tabular analysis completed for {len(outputs)} file(s)."}


def _fallback_summarize(paths: list[str], output_mode: str, artifacts_dir: str) -> dict:
    doc_paths = [p for p in paths if os.path.isfile(p)]
    if not doc_paths:
        return {"success": False, "summary": "No readable document path provided for fallback summary."}
    outputs: list[str] = []
    for doc_path in doc_paths:
        extracted = _extract_document_text(doc_path)
        bullets = _make_bullet_points(extracted, max_points=5)
        words = len(extracted.split()) if extracted else 0
        size = os.path.getsize(doc_path) if os.path.exists(doc_path) else 0
        lines = [
            f"# Summary: {os.path.basename(doc_path)}",
            "",
            "## Document Info",
            f"- Type: {os.path.splitext(doc_path)[1].lower() or 'unknown'}",
            f"- Size: {size} bytes",
            f"- Extracted words: {words}",
            "",
            "## Key Points",
        ]
        if bullets:
            lines.extend([f"- {b}" for b in bullets])
        else:
            lines.append("- No text extracted; binary or unsupported format.")
        lines.extend([
            "",
            "## Executive Summary",
            "Deterministic fallback summary generated after agentic runtime failure.",
            "",
        ])
        out_path = _target_output_path(output_mode, doc_path, artifacts_dir, ".summary.md")
        _write_text(out_path, "\n".join(lines))
        outputs.append(out_path)
    return {"success": True, "summary": f"Document summary completed for {len(outputs)} file(s)."}


def _fallback_organize(paths: list[str], output_mode: str, artifacts_dir: str) -> dict:
    folders = [p for p in paths if os.path.isdir(p)]
    if not folders:
        return {"success": False, "summary": "No folder path provided for fallback organization."}
    outputs: list[str] = []
    ext_groups = {
        "documents": {".pdf", ".doc", ".docx", ".docm", ".dotx", ".dotm"},
        "text": {".txt", ".md", ".rtf"},
        "data": {".csv", ".xlsx", ".xls"},
        "images": {".png", ".jpg", ".jpeg", ".gif", ".svg"},
        "presentations": {".ppt", ".pptx"},
        "code": {".py", ".js", ".ts", ".java", ".c", ".cpp", ".h"},
    }
    for folder in folders:
        by_category: dict[str, list[str]] = {}
        by_date_folder: dict[str, list[str]] = {}
        for root, dirs, files in os.walk(folder):
            dirs[:] = [d for d in dirs if not d.startswith(".")]
            rel_root = os.path.relpath(root, folder)
            date_key = ""
            if rel_root != ".":
                parts = rel_root.split(os.sep)
                for part in parts:
                    if re.fullmatch(r"\d{8}", part) or re.fullmatch(r"\d{4}-\d{2}-\d{2}", part):
                        date_key = part
                        break
            for name in files:
                if name.startswith("."):
                    continue
                ext = os.path.splitext(name)[1].lower()
                rel = os.path.relpath(os.path.join(root, name), folder)
                category = "other"
                for cat, exts in ext_groups.items():
                    if ext in exts:
                        category = cat
                        break
                by_category.setdefault(category, []).append(rel)
                if date_key:
                    by_date_folder.setdefault(date_key, []).append(rel)
        lines = [
            "# Organization Plan",
            "",
            "## Source Folder",
            f"- {folder}",
            "",
            "## Proposed Grouping by Category",
        ]
        for cat in sorted(by_category):
            lines.append(f"### {cat.title()} ({len(by_category[cat])})")
            for rel in sorted(by_category[cat])[:50]:
                lines.append(f"- {rel}")
            if len(by_category[cat]) > 50:
                lines.append(f"- ... {len(by_category[cat]) - 50} more")
            lines.append("")
        lines.append("## Date-Based Grouping (if detected)")
        if by_date_folder:
            for key in sorted(by_date_folder):
                lines.append(f"### {key} ({len(by_date_folder[key])})")
                for rel in sorted(by_date_folder[key])[:50]:
                    lines.append(f"- {rel}")
                if len(by_date_folder[key]) > 50:
                    lines.append(f"- ... {len(by_date_folder[key]) - 50} more")
                lines.append("")
        else:
            lines.append("- No date-based folder structure detected.")
        lines.append("## Notes")
        lines.append("- Plan generated by deterministic fallback after agentic runtime failure.")
        lines.append("- No files were moved or deleted.")
        filename = "organization-plan.md"
        if output_mode == "workspace" and len(folders) > 1:
            filename = f"{os.path.basename(folder.rstrip('/'))}.organization-plan.md"
        out_path = _target_output_file(output_mode, folder, artifacts_dir, filename)
        _write_text(out_path, "\n".join(lines) + "\n")
        outputs.append(out_path)
    return {"success": True, "summary": f"Organization plan generated for {len(outputs)} folder(s)."}


def _run_deterministic_fallback(capability: str, paths: list[str], output_mode: str, artifacts_dir: str) -> dict:
    if capability == "analyze":
        return _fallback_analyze(paths, output_mode, artifacts_dir)
    if capability == "summarize":
        return _fallback_summarize(paths, output_mode, artifacts_dir)
    if capability == "organize":
        return _fallback_organize(paths, output_mode, artifacts_dir)
    return {"success": False, "summary": f"Fallback unsupported capability: {capability}"}


def _build_summarize_prompt(paths: list[str], output_mode: str, source_root: str) -> str:
    paths_list = "\n".join(f"- {p}" for p in paths)
    has_multiple_files = len(paths) > 1
    target_lines = []
    for path in paths:
        if output_mode == "workspace":
            target_lines.append(f"- Source: {path}\n  Target filename: {os.path.basename(path)}.summary.md")
        else:
            target_lines.append(f"- Source: {path}\n  Target path: {path}.summary.md")
    if has_multiple_files:
        if output_mode == "workspace":
            target_lines.append("- Combined report target filename: combined-summary.md")
        else:
            target_lines.append(f"- Combined report target path: {os.path.join(os.path.dirname(paths[0]), 'combined-summary.md')}")
    targets_block = "\n".join(target_lines)
    write_rules = (
        "2. Write a summary using the write_workspace tool to the exact target filename listed below."
        if output_mode == "workspace"
        else
        "2. Write a summary using the write_file tool to the exact target path listed below."
    )
    return f"""Summarize the following document(s):

{paths_list}

Source root: {source_root}
Output mode: {output_mode}
Required output targets:
{targets_block}

For each file:
1. Read the file using the appropriate tool:
   - PDF: `read_pdf`
   - Word OpenXML (`.docx/.docm/.dotx/.dotm`): `read_docx`
   - Plain text / Markdown / HTML / XML / TSV: `read_txt`
   - CSV: `read_csv`
   - Excel (`.xlsx/.xls`): `read_xlsx` or `read_xls`
   - PowerPoint (`.pptx`): `read_pptx`
{write_rules}
3. Summarize in English. For French or other foreign language documents, summarize accurately.
4. Use only the provided MCP tools. Do not use native Write/Edit/Read tools.
5. Preserve the full original filename, including its extension, before appending `.summary.md`.
6. If multiple documents are provided, you MUST create `combined-summary.md` that includes one section per document and a cross-document overview.

Output format for each file:
# Summary: {{filename}}

## Document Info
- Type: PDF/Word/TXT/Markdown/HTML/XML/PPTX/Spreadsheet
- Size: N KB

## Key Points
- Point 1
- Point 2
- Point 3

## Executive Summary
1-paragraph summary of the document.

If a PDF or document has no extractable embedded text with the provided tools, state that explicitly in English and do not invent content.

If multiple files are provided, `combined-summary.md` is mandatory.
"""


def _build_analyze_prompt(paths: list[str], output_mode: str, source_root: str) -> str:
    paths_list = "\n".join(f"- {p}" for p in paths)
    target_lines = []
    for path in paths:
        if output_mode == "workspace":
            target_lines.append(f"- Source: {path}\n  Target filename: {os.path.basename(path)}.analysis.md")
        else:
            target_lines.append(f"- Source: {path}\n  Target path: {path}.analysis.md")
    targets_block = "\n".join(target_lines)
    write_rules = (
        "4. Write an analysis report using write_workspace to the exact target filename listed below."
        if output_mode == "workspace"
        else
        "4. Write an analysis report using write_file to the exact target path listed below."
    )
    return f"""Analyze the following tabular/document data source(s):

{paths_list}

Source root: {source_root}
Output mode: {output_mode}
Required output targets:
{targets_block}

Methodology (must follow):
1. Inspect each source first to infer schema and structure:
   - CSV/TSV: read_csv or read_txt (if delimiter-based text)
   - XLSX: read_xlsx
   - XLS: read_xls
   - Document-based tables: read_pdf/read_docx/read_txt as needed
2. Explicitly state inferred schema: columns/fields, inferred data types, missing-value patterns.
3. Produce analysis from inferred schema only (no hardcoded column-name assumptions):
   - Summary statistics for detected numeric fields
   - Relevant aggregations for categorical fields
   - Trends, anomalies, and data-quality caveats
{write_rules}
5. Use only the provided MCP tools. Do not use native Write/Edit/Read tools.

Output format:
# Data Analysis: {{filename_or_source}}

## File Overview
- Source type and parse method
- Rows/records (if structured)
- Fields/columns detected

## Summary Statistics
Use Markdown tables where possible. For each detected numeric field include: count, min, max, average.

## Key Insights
- Schema-driven insights only (no fixed business template)
- Mention assumptions and confidence limits

IMPORTANT: Write all analysis results in English. Do not use any other language.

CRITICAL:
- Preserve the full original filename, including its extension, before appending `.analysis.md`.
- Do not write to relative paths like `artifacts/...`. The only authorized workspace write path is via write_workspace.
"""


def _build_organize_prompt(paths: list[str], output_mode: str, source_root: str) -> str:
    paths_list = "\n".join(f"- {p}" for p in paths)
    write_rules = (
        "3. Write the organization plan using write_workspace tool with filename: organization-plan.md"
        if output_mode == "workspace"
        else
        "3. Write the organization plan using write_file tool to: {source_folder}/organization-plan.md"
    )
    return f"""Analyze and organize the following folder(s):

{paths_list}

Source root: {source_root}
Output mode: {output_mode}

TASK:
Your goal is to discover meaningful patterns in the file content and structure, then CREATE the organized folder structure in the workspace.

WORKFLOW:
1. Call organize_folder on each source folder first. Treat its recursive `files` inventory as the authoritative source list.
2. Review the returned per-file metadata:
   - Use `primary_entity`, `primary_entity_source`, `primary_entity_confidence`, `inferred_date_bucket`, `prominent_headings`, `labeled_fields`, `suggested_reader_tool`, and `suggested_destination`
   - Use `entity_counts` and `date_bucket_counts` to sanity-check the overall distribution before copying
   - If `primary_entity_confidence` is `high`, treat that identity as authoritative
   - If metadata is ambiguous or missing, inspect representative files with the suggested reader tool
   - Do NOT infer ownership from assignment titles, book titles, or topic headings
3. Based on the discovered patterns, determine the BEST grouping strategy:
   - Choose grouping criteria that meaningfully organizes the files
   - Examples: by type, by date, by author, by status, by project, by topic, etc.
4. EXECUTE the organization using organize_move_file tool:
   - Use mkdir action to create directory structure under `organized-output/files/`
   - Use copy_file action to copy files to their organized locations under `organized-output/files/`
   - In workspace mode, pass category-relative destinations such as:
     `entities/Entity_A/YYYY-MM/source-001.txt`
     `entities/Entity_B/YYYY-MM/source-014.txt`
     The tool will place them under `organized-output/files/`
   - Write the organization plan using write_workspace tool with filename: organization-plan.md

CRITICAL: You must USE the organize_move_file tool to actually create the organized folder structure. Do not just write a plan - execute it.
CRITICAL: A plan-only answer is a failure. The task is complete only if files exist under `organized-output/files/`.
CRITICAL: Every non-hidden source file from `organize_folder.files` must be copied exactly once. Do not duplicate a source file into multiple destinations.
CRITICAL: When `primary_entity` or `inferred_date_bucket` is present in the organize_folder metadata, use it in the destination path unless a file inspection proves it is wrong.
CRITICAL: Never reference or copy a source path that is not present in `organize_folder.files`.
CRITICAL: If `suggested_destination` is present for a file, use that exact relative destination unless a direct file inspection proves it is wrong.
CRITICAL: Your final response and the contents of `organization-plan.md` must be in English only.

OUTPUT FORMAT:
Write a summary to organization-plan.md explaining:
# Folder Organization Plan

## Discovered Patterns
What patterns you found in the files.

## Organized Structure Created
Show the actual directory structure you created under `organized-output/files/`.

## Files Organized
List which files were moved to which locations.
"""


# ---------------------------------------------------------------------------
# Node: report_result
# ---------------------------------------------------------------------------

def report_result(state: dict) -> dict:
    """Write pr-evidence.json, warnings.md (if partial failures), and return final result."""
    import time
    workspace_root = state.get("workspace_root", "")
    capability = state.get("capability", "summarize")
    summary = state.get("summary", "")
    validated_paths = state.get("validated_paths", [])
    output_mode = state.get("output_mode", "workspace")
    warnings = state.get("warnings", [])
    success = bool(state.get("success", False))
    status = "completed" if success else "failed"
    expected_outputs = state.get("expected_outputs", [])
    raw_output = state.get("raw_output", "")
    english_summary = summary if summary and not _contains_cjk(summary) else _english_summary_for_report(capability, success, expected_outputs)

    if workspace_root:
        os.makedirs(workspace_root, exist_ok=True)

        # Write warnings.md if there are warnings
        if warnings:
            warnings_path = os.path.join(workspace_root, "warnings.md")
            try:
                with open(warnings_path, "w", encoding="utf-8") as f:
                    f.write("# Warnings\n\n")
                    for w in warnings:
                        f.write(f"- {w}\n")
                logger.info(f"report_result: warnings written to {warnings_path}")
            except OSError as exc:
                logger.error(f"report_result: failed to write warnings: {exc}")

        # Write task-report.json (was pr-evidence.json)
        task_report_path = os.path.join(workspace_root, "task-report.json")
        evidence = {
            "metadata": {
                "agent_id": AGENT_ID,
                "step": "report_result",
                "timestamp": time.strftime("%Y-%m-%dT%H:%M:%S%z"),
            },
            "data": {
                "capability": capability,
                "output_mode": output_mode,
                "source_paths": validated_paths,
                "summary": english_summary[:500] if english_summary else "",
                "success": success,
                "artifacts_dir": state.get("artifacts_dir", ""),
                "expected_outputs": expected_outputs,
                "warnings_count": len(warnings),
            },
        }
        try:
            with open(task_report_path, "w", encoding="utf-8") as fh:
                json.dump(evidence, fh, ensure_ascii=False, indent=2)
            logger.info(f"report_result: task-report written to {task_report_path}")
        except OSError as exc:
            logger.error(f"report_result: failed to write task-report: {exc}")

        if raw_output:
            raw_output_path = os.path.join(workspace_root, "agentic-output.txt")
            try:
                with open(raw_output_path, "w", encoding="utf-8") as fh:
                    fh.write(_english_agentic_output(raw_output, capability, expected_outputs))
                logger.info(f"report_result: agentic output written to {raw_output_path}")
            except OSError as exc:
                logger.error(f"report_result: failed to write agentic output: {exc}")

    return {
        "status": status,
        "summary": english_summary,
        "capability": capability,
        "output_mode": output_mode,
        "success": success,
        "warnings_count": len(warnings),
    }
