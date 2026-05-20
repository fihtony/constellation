import pytest, os, tempfile, json

def test_read_pptx_tool_with_real_file():
    """Test ReadPptxTool with a real PPTX file from tests/data/."""
    from agents.office.office_tools import ReadPptxTool

    # Find a real pptx file in tests/data
    import glob
    pptx_files = glob.glob("tests/data/**/*.pptx", recursive=True)

    if not pptx_files:
        # Create a minimal test pptx using python-pptx
        import tempfile
        from pptx import Presentation

        try:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            if title:
                title.text = "Test Slide Title"
            body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if body:
                body.text = "Test body content"

            tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
            tmp_path = tmp.name
            tmp.close()
            prs.save(tmp_path)
            pptx_files = [tmp_path]
        except Exception:
            pytest.skip("No .pptx files found and could not create test file")

    tool = ReadPptxTool()
    old_root = os.environ.get("OFFICE_SOURCE_ROOT")
    try:
        # Use the parent dir of the pptx file as source root so path validation passes
        source_root = os.path.dirname(os.path.abspath(pptx_files[0]))
        os.environ["OFFICE_SOURCE_ROOT"] = source_root
        result = tool.execute_sync(path=pptx_files[0])
        assert result.success, f"read_pptx failed: {result.error}"
        data = json.loads(result.output)
        assert "content" in data
        assert "slides" in data
        assert "total_slides" in data
    finally:
        if old_root is not None:
            os.environ["OFFICE_SOURCE_ROOT"] = old_root
        else:
            os.environ.pop("OFFICE_SOURCE_ROOT", None)

    # Clean up temp file if we created one
    if pptx_files and pptx_files[0].startswith("/var/folders"):
        try:
            os.unlink(pptx_files[0])
        except Exception:
            pass


def test_read_pptx_rejects_ppt_format():
    """Test that ReadPptxTool rejects .ppt files with clear error."""
    from agents.office.office_tools import ReadPptxTool

    tool = ReadPptxTool()
    # Create a fake .ppt file inside tests/data so it passes path validation
    pptx_file = "tests/data/fake.ppt"
    old_root = os.environ.get("OFFICE_SOURCE_ROOT")
    try:
        os.environ["OFFICE_SOURCE_ROOT"] = os.path.abspath("tests/data")
        # Write a minimal file (not real pptx but has .ppt extension)
        with open(pptx_file, "wb") as f:
            f.write(b"fake ppt content")
        result = tool.execute_sync(path=pptx_file)
        assert not result.success
        assert "convert" in result.error.lower() and ".ppt" in result.error.lower()
    finally:
        if old_root is not None:
            os.environ["OFFICE_SOURCE_ROOT"] = old_root
        else:
            os.environ.pop("OFFICE_SOURCE_ROOT", None)
        if os.path.exists(pptx_file):
            os.unlink(pptx_file)


def test_read_pptx_rejects_outside_source_root():
    """Test path validation blocks access outside OFFICE_SOURCE_ROOT."""
    from agents.office.office_tools import ReadPptxTool

    tool = ReadPptxTool()
    old_root = os.environ.get("OFFICE_SOURCE_ROOT")
    try:
        os.environ["OFFICE_SOURCE_ROOT"] = "/tmp"
        result = tool.execute_sync(path="/etc/some.pptx")
        assert not result.success
        assert "outside OFFICE_SOURCE_ROOT" in result.error
    finally:
        if old_root is not None:
            os.environ["OFFICE_SOURCE_ROOT"] = old_root
        else:
            os.environ.pop("OFFICE_SOURCE_ROOT", None)