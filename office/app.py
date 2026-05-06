"""Office Agent — local document summary, analysis, and organize execution agent."""

from __future__ import annotations

import csv
import json
import os
import re
import shutil
import threading
import time
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from urllib.parse import urlparse
from urllib.request import Request, urlopen

from common.instance_reporter import InstanceReporter
from common.message_utils import build_text_artifact, extract_text
from common.orchestrator import resolve_orchestrator_base_url
from common.per_task_exit import PerTaskExitHandler
from common.tools.control_tools import configure_control_tools
from common.rules_loader import build_system_prompt
from common.prompt_builder import build_system_prompt_from_manifest
from common.agent_system_prompt import build_agent_system_prompt as _build_manifest_prompt
from common.runtime.adapter import get_runtime, require_agentic_runtime, summarize_runtime_configuration
from common.task_permissions import (
    PermissionDeniedError,
    audit_permission_check,
    build_permission_denied_artifact,
    build_permission_denied_details,
    parse_permission_grant,
)
from common.task_store import TaskStore
from common.time_utils import local_iso_timestamp
from office import prompts

from common.env_utils import load_dotenv

load_dotenv(os.path.join(os.path.dirname(__file__), ".env"))

HOST = os.environ.get("HOST", "0.0.0.0")
PORT = int(os.environ.get("PORT", "8060"))
AGENT_ID = os.environ.get("AGENT_ID", "office-agent")
ADVERTISED_URL = os.environ.get("ADVERTISED_BASE_URL", f"http://office-agent:{PORT}")

ACK_TIMEOUT = int(os.environ.get("A2A_ACK_TIMEOUT_SECONDS", "15"))
TASK_TIMEOUT = int(os.environ.get("A2A_TASK_TIMEOUT_SECONDS", "1200"))
MAX_FILE_SIZE_BYTES = int(os.environ.get("OFFICE_MAX_FILE_SIZE_MB", "50")) * 1024 * 1024
MAX_DIR_FILE_COUNT = int(os.environ.get("OFFICE_MAX_DIR_FILE_COUNT", "2000"))
MAX_DIR_TOTAL_BYTES = int(os.environ.get("OFFICE_MAX_DIR_TOTAL_MB", "250")) * 1024 * 1024
INPUT_ROOT = "/app/userdata"

task_store = TaskStore()
exit_handler = PerTaskExitHandler()
reporter = InstanceReporter(
    agent_id=AGENT_ID,
    service_url=ADVERTISED_URL,
    port=PORT,
)

_SERVER: ThreadingHTTPServer | None = None
_WRITE_ROOTS_LOCK = threading.Lock()
_ACTIVE_WRITE_ROOTS: set[str] = set()
_ORGANIZE_SCHEMA_ROOT = "files"
_ORGANIZE_WRAPPER_SEGMENTS = {
    "organized-output",
    "grouped",
    "by-student",
    "originals",
    "output",
    "outputs",
    "results",
    "workspace",
    "final",
}


def audit_log(event: str, **kwargs):
    entry = {"ts": local_iso_timestamp(), "event": event, **kwargs}
    print(f"[audit] {json.dumps(entry, ensure_ascii=False)}")


_MANIFEST_SYSTEM_PROMPT: str = ""


def _get_manifest_system_prompt() -> str:
    """Return cached manifest-based system prompt, building it on first call."""
    global _MANIFEST_SYSTEM_PROMPT
    if not _MANIFEST_SYSTEM_PROMPT:
        agent_dir = os.path.dirname(os.path.abspath(__file__))
        _MANIFEST_SYSTEM_PROMPT = build_system_prompt_from_manifest(agent_dir) or build_system_prompt(
            prompts.SYSTEM, "office"
        )
    return _MANIFEST_SYSTEM_PROMPT


def _runtime_config_summary() -> dict:
    return {
        "service": AGENT_ID,
        "maxFileSizeBytes": MAX_FILE_SIZE_BYTES,
        "maxDirFileCount": MAX_DIR_FILE_COUNT,
        "maxDirTotalBytes": MAX_DIR_TOTAL_BYTES,
        "runtimeConfig": summarize_runtime_configuration(),
    }


def _permission_enforcement_mode() -> str:
    return os.environ.get("PERMISSION_ENFORCEMENT", "strict").strip().lower() or "strict"


def _check_office_permission(
    *,
    action: str,
    target: str,
    metadata: dict,
    scope: str = "*",
) -> tuple[bool, str, str]:
    if _permission_enforcement_mode() == "off":
        return True, "allowed", ""

    request_agent = str(metadata.get("requestAgent") or "compass-agent").strip() or "compass-agent"
    task_id = str(metadata.get("orchestratorTaskId") or metadata.get("taskId") or "").strip()
    permissions_data = metadata.get("permissions") if isinstance(metadata.get("permissions"), dict) else None
    grant = parse_permission_grant(permissions_data)
    if grant:
        allowed, reason = grant.check("office", action, scope)
        escalation = grant.escalation_for("office", action, scope)
    else:
        allowed = False
        reason = "No permissions attached to request. Explicit permission grant required."
        escalation = "require_user_approval"

    audit_permission_check(
        task_id=task_id,
        orchestrator_task_id=task_id,
        request_agent=request_agent,
        target_agent=AGENT_ID,
        action=action,
        target=target,
        decision="allowed" if allowed else "denied",
        reason=reason,
        agent_id=AGENT_ID,
    )
    return allowed, reason, escalation


def _require_office_permission(
    *,
    action: str,
    target: str,
    metadata: dict,
    scope: str = "*",
) -> None:
    allowed, reason, escalation = _check_office_permission(
        action=action,
        target=target,
        metadata=metadata,
        scope=scope,
    )
    if allowed:
        return
    if _permission_enforcement_mode() == "strict":
        raise PermissionDeniedError(
            build_permission_denied_details(
                permission_agent="office",
                target_agent=AGENT_ID,
                action=action,
                target=target,
                reason=reason,
                escalation=escalation or "require_user_approval",
                scope=scope,
                request_agent=str(metadata.get("requestAgent") or "compass-agent").strip() or "compass-agent",
                task_id=str(metadata.get("taskId") or ""),
                orchestrator_task_id=str(metadata.get("orchestratorTaskId") or ""),
            )
        )

    print(f"[{AGENT_ID}] WARN: permission check failed but enforcement={_permission_enforcement_mode()}: {reason}")


def _enforce_office_task_permissions(metadata: dict, target_paths: list[str], input_root: str, output_mode: str) -> None:
    for path in target_paths:
        if not _path_within_base(path, input_root):
            _require_office_permission(
                action="access_outside_root",
                target=path,
                metadata=metadata,
                scope="*",
            )
            raise RuntimeError(f"Target path escapes mounted input root: {path}")

    root_target = os.path.commonpath(target_paths)
    _require_office_permission(
        action="read",
        target=root_target,
        metadata=metadata,
        scope="task_root",
    )
    if output_mode == "inplace":
        _require_office_permission(
            action="write",
            target=root_target,
            metadata=metadata,
            scope="task_root",
        )


def _parse_json(text: str) -> dict:
    text = (text or "").strip()
    if not text:
        return {}
    if text.startswith("```"):
        lines = text.splitlines()
        start = 1
        end = len(lines)
        while end > start and lines[end - 1].strip() in ("```", ""):
            end -= 1
        text = "\n".join(lines[start:end]).strip()
    try:
        payload = json.loads(text)
        return payload if isinstance(payload, dict) else {}
    except json.JSONDecodeError:
        pass
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if not match:
        return {}
    try:
        payload = json.loads(match.group())
        return payload if isinstance(payload, dict) else {}
    except json.JSONDecodeError:
        return {}


def _run_agentic_json(
    prompt: str,
    actor: str,
    *,
    system_prompt: str,
    context: dict | None = None,
    timeout: int = 180,
) -> dict:
    require_agentic_runtime("Office Agent")
    result = get_runtime().run(
        prompt=prompt,
        context=context,
        system_prompt=system_prompt,
        timeout=timeout,
        max_tokens=4096,
    )
    for warning in result.get("warnings") or []:
        print(f"[{AGENT_ID}] Runtime warning ({actor}): {warning}")
    return _parse_json(result.get("raw_response") or result.get("summary") or "")


def _path_within_base(path: str, base: str) -> bool:
    try:
        common = os.path.commonpath([os.path.realpath(path), os.path.realpath(base)])
    except ValueError:
        return False
    return common == os.path.realpath(base)


def _safe_output_path(root: str, relative_path: str) -> str:
    relative = str(relative_path or "").strip().replace("\\", "/")
    if not relative or os.path.isabs(relative):
        raise RuntimeError(f"Destination must be a non-empty relative path: {relative_path}")
    candidate = os.path.realpath(os.path.join(root, relative))
    if not _path_within_base(candidate, root):
        raise RuntimeError(f"Destination escapes the allowed output root: {relative_path}")
    return candidate


def _workspace_root(metadata: dict) -> str:
    workspace = str(metadata.get("officeWorkspacePath") or metadata.get("sharedWorkspacePath") or "").strip()
    if not workspace:
        raise RuntimeError("Office task is missing a shared workspace path.")
    os.makedirs(workspace, exist_ok=True)
    return workspace


def _audit_dir(metadata: dict) -> str:
    root = os.path.join(_workspace_root(metadata), AGENT_ID)
    os.makedirs(root, exist_ok=True)
    return root


def _write_text_file(path: str, content: str) -> None:
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "w", encoding="utf-8") as handle:
        handle.write(content)


def _save_audit_file(metadata: dict, relative_name: str, content: str) -> str:
    full_path = os.path.join(_audit_dir(metadata), relative_name)
    _write_text_file(full_path, content)
    return full_path


def _append_command_log(metadata: dict, line: str) -> None:
    path = os.path.join(_audit_dir(metadata), "command-log.txt")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "a", encoding="utf-8") as handle:
        handle.write(f"[{local_iso_timestamp()}] {line}\n")


def _update_stage_summary(metadata: dict, task_id: str, phase: str, **extra) -> None:
    payload = {
        "taskId": task_id,
        "currentPhase": phase,
        "updatedAt": local_iso_timestamp(),
        "runtimeConfig": _runtime_config_summary(),
    }
    payload.update({key: value for key, value in extra.items() if value is not None})
    _save_audit_file(metadata, "stage-summary.json", json.dumps(payload, ensure_ascii=False, indent=2))


def _notify_callback(callback_url: str, task_id: str, state: str, status_message: str, artifacts: list | None = None):
    if not callback_url:
        return
    payload = {
        "downstreamTaskId": task_id,
        "state": state,
        "statusMessage": status_message,
        "artifacts": artifacts or [],
        "agentId": AGENT_ID,
    }
    data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    request = Request(
        callback_url,
        data=data,
        headers={"Content-Type": "application/json; charset=utf-8"},
        method="POST",
    )
    try:
        with urlopen(request, timeout=10):
            pass
    except Exception as err:
        print(f"[{AGENT_ID}] Callback failed: {err}")


def _report_progress(compass_url: str, compass_task_id: str, step: str) -> None:
    if not compass_url or not compass_task_id or not step:
        return
    payload = {"step": step, "agentId": AGENT_ID}
    data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    request = Request(
        f"{compass_url.rstrip('/')}/tasks/{compass_task_id}/progress",
        data=data,
        headers={"Content-Type": "application/json; charset=utf-8"},
        method="POST",
    )
    try:
        with urlopen(request, timeout=5):
            pass
    except Exception as err:
        print(f"[{AGENT_ID}] Progress report failed (non-critical): {err}")


def _schedule_shutdown(delay_seconds: int = 5):
    def _shutdown():
        time.sleep(delay_seconds)
        print(f"[{AGENT_ID}] Per-task shutdown triggered")
        if _SERVER:
            _SERVER.shutdown()

    threading.Thread(target=_shutdown, daemon=True).start()


def _auto_stop_after_task_enabled() -> bool:
    return os.environ.get("AUTO_STOP_AFTER_TASK", "").strip() == "1"


def _apply_task_exit_rule(task_id: str, exit_rule: dict) -> None:
    def _run():
        rule_type = (exit_rule or {}).get("type", "wait_for_parent_ack")
        if rule_type == "auto_stop":
            if not _auto_stop_after_task_enabled():
                print(f"[{AGENT_ID}] AUTO_STOP_AFTER_TASK not set — keeping agent alive")
                return
            rule_type = "immediate"
        exit_handler.apply(
            task_id,
            {**(exit_rule or {}), "type": rule_type},
            shutdown_fn=_schedule_shutdown,
            agent_id=AGENT_ID,
        )

    threading.Thread(target=_run, daemon=True).start()


def _detect_encoding(path: str) -> str:
    try:
        import chardet  # type: ignore
    except Exception:
        return "utf-8"
    with open(path, "rb") as handle:
        raw = handle.read(8192)
    detected = chardet.detect(raw or b"")
    return detected.get("encoding") or "utf-8"


def _read_txt(path: str) -> str:
    encoding = _detect_encoding(path)
    with open(path, "r", encoding=encoding, errors="replace") as handle:
        return handle.read()


def _read_csv_rows(path: str, limit: int | None = 100) -> tuple[list[dict], list[str]]:
    encoding = _detect_encoding(path)
    with open(path, "r", encoding=encoding, newline="", errors="replace") as handle:
        reader = csv.DictReader(handle)
        rows = []
        for row in reader:
            rows.append(dict(row))
            if limit is not None and len(rows) >= limit:
                break
        return rows, list(reader.fieldnames or [])


def _coerce_float(value) -> float | None:
    text = str(value or "").strip()
    if not text:
        return None
    text = text.replace(",", "").replace("$", "")
    try:
        return float(text)
    except (TypeError, ValueError):
        return None


def _build_grouped_numeric_totals(rows: list[dict], fields: list[str], numeric_fields: list[str], max_groups: int = 20) -> dict:
    grouped_totals: dict[str, dict[str, list[dict]]] = {}
    for field in fields:
        if field in numeric_fields:
            continue
        unique_values = []
        seen = set()
        for row in rows:
            value = str(row.get(field) or "").strip()
            if not value or value in seen:
                continue
            seen.add(value)
            unique_values.append(value)
            if len(unique_values) > max_groups:
                break
        if len(unique_values) < 2 or len(unique_values) > max_groups:
            continue

        field_totals: dict[str, list[dict]] = {}
        for numeric_field in numeric_fields:
            buckets: dict[str, dict[str, float]] = {}
            for row in rows:
                group = str(row.get(field) or "").strip()
                numeric_value = _coerce_float(row.get(numeric_field))
                if not group or numeric_value is None:
                    continue
                bucket = buckets.setdefault(group, {"sum": 0.0, "count": 0.0})
                bucket["sum"] += numeric_value
                bucket["count"] += 1.0
            if not buckets:
                continue
            ranked = sorted(
                (
                    {
                        "group": group,
                        "sum": round(values["sum"], 4),
                        "count": int(values["count"]),
                        "avg": round(values["sum"] / values["count"], 4),
                    }
                    for group, values in buckets.items()
                ),
                key=lambda item: (-item["sum"], item["group"]),
            )
            field_totals[numeric_field] = ranked[:max_groups]
        if field_totals:
            grouped_totals[field] = field_totals
    return grouped_totals


def _read_excel_workbook(path: str, limit_rows: int = 100) -> dict[str, list[list[object]]]:
    suffix = os.path.splitext(path)[1].lower()
    if suffix == ".xlsx":
        try:
            import openpyxl  # type: ignore
        except Exception as exc:
            raise RuntimeError(f"Missing openpyxl dependency for .xlsx: {exc}") from exc
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
        result = {}
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            rows = []
            for row_index, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                rows.append(list(row))
                if row_index >= limit_rows:
                    break
            result[sheet_name] = rows
        return result

    try:
        import xlrd  # type: ignore
    except Exception as exc:
        raise RuntimeError(f"Missing xlrd dependency for .xls: {exc}") from exc
    workbook = xlrd.open_workbook(path)
    result = {}
    for sheet in workbook.sheets():
        rows = []
        for row_index in range(min(sheet.nrows, limit_rows)):
            rows.append(sheet.row_values(row_index))
        result[sheet.name] = rows
    return result


def _read_docx(path: str) -> str:
    try:
        from docx import Document  # type: ignore
    except Exception as exc:
        raise RuntimeError(f"Missing python-docx dependency: {exc}") from exc
    document = Document(path)
    return "\n".join(para.text for para in document.paragraphs if para.text.strip())


def _read_pptx(path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        raise RuntimeError(f"Missing python-pptx dependency: {exc}") from exc
    presentation = Presentation(path)
    texts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and str(shape.text).strip():
                texts.append(str(shape.text))
    return "\n".join(texts)


def _read_pdf(path: str) -> str:
    try:
        import pdfplumber  # type: ignore
    except Exception as exc:
        raise RuntimeError(f"Missing pdfplumber dependency: {exc}") from exc
    texts = []
    with pdfplumber.open(path) as document:
        for page in document.pages:
            text = page.extract_text() or ""
            if text.strip():
                texts.append(text)
    if not texts:
        raise RuntimeError("No extractable text found in PDF. OCR/scanned PDFs are not supported.")
    return "\n".join(texts)


def _extract_document_preview(path: str) -> dict:
    suffix = os.path.splitext(path)[1].lower()
    if suffix in {".doc", ".ppt"}:
        raise RuntimeError(f"Legacy Office format {suffix} is not supported in MVP; please convert it first.")

    if suffix == ".txt":
        text = _read_txt(path)
        return {"path": path, "type": suffix, "preview": text[:4000]}
    if suffix == ".csv":
        rows, fields = _read_csv_rows(path)
        preview = json.dumps(rows[:10], ensure_ascii=False, indent=2)
        return {"path": path, "type": suffix, "preview": preview, "fields": fields}
    if suffix in {".xlsx", ".xls"}:
        workbook = _read_excel_workbook(path)
        return {
            "path": path,
            "type": suffix,
            "preview": json.dumps(workbook, ensure_ascii=False, indent=2)[:4000],
            "sheets": list(workbook.keys()),
        }
    if suffix == ".docx":
        return {"path": path, "type": suffix, "preview": _read_docx(path)[:4000]}
    if suffix == ".pptx":
        return {"path": path, "type": suffix, "preview": _read_pptx(path)[:4000]}
    if suffix == ".pdf":
        return {"path": path, "type": suffix, "preview": _read_pdf(path)[:4000]}

    raise RuntimeError(f"Unsupported file type: {suffix}")


def _preflight_scan(target_paths: list[str]) -> dict:
    """Quick size/count scan without reading file content. Returns stats and over-limit flags."""
    file_count = 0
    total_bytes = 0
    large_files: list[dict] = []
    for target in target_paths:
        real = os.path.realpath(target)
        if os.path.isfile(real):
            sz = os.path.getsize(real)
            file_count += 1
            total_bytes += sz
            if sz > MAX_FILE_SIZE_BYTES:
                large_files.append({"path": target, "sizeBytes": sz})
        elif os.path.isdir(real):
            for root, _, names in os.walk(real):
                for name in names:
                    child = os.path.join(root, name)
                    if os.path.isfile(child):
                        sz = os.path.getsize(child)
                        file_count += 1
                        total_bytes += sz
                        if sz > MAX_FILE_SIZE_BYTES:
                            large_files.append({"path": child, "sizeBytes": sz})
    return {
        "fileCount": file_count,
        "totalBytes": total_bytes,
        "largeFiles": large_files[:10],
        "overFileCountLimit": file_count > MAX_DIR_FILE_COUNT,
        "overBytesLimit": total_bytes > MAX_DIR_TOTAL_BYTES,
        "limitFileCount": MAX_DIR_FILE_COUNT,
        "limitTotalMB": MAX_DIR_TOTAL_BYTES // (1024 * 1024),
    }


def _collect_files(target_paths: list[str], *, allow_any: bool = False) -> tuple[list[str], list[str]]:
    files: list[str] = []
    warnings: list[str] = []
    total_bytes = 0

    def _include_file(path: str):
        nonlocal total_bytes
        size_bytes = os.path.getsize(path)
        if size_bytes > MAX_FILE_SIZE_BYTES:
            warnings.append(f"Skipped oversized file (>{MAX_FILE_SIZE_BYTES // (1024 * 1024)} MB): {path}")
            return
        total_bytes += size_bytes
        files.append(path)

    for target in target_paths:
        real_target = os.path.realpath(target)
        if os.path.isfile(real_target):
            size_bytes = os.path.getsize(real_target)
            if size_bytes > MAX_FILE_SIZE_BYTES:
                raise RuntimeError(
                    f"File exceeds size limit ({MAX_FILE_SIZE_BYTES // (1024 * 1024)} MB): {real_target}"
                )
            total_bytes += size_bytes
            files.append(real_target)
            continue

        count = 0
        for root, _, names in os.walk(real_target):
            for name in names:
                child = os.path.realpath(os.path.join(root, name))
                count += 1
                if count > MAX_DIR_FILE_COUNT:
                    raise RuntimeError(
                        f"Directory exceeds file count limit ({MAX_DIR_FILE_COUNT} files): {real_target}"
                    )
                if allow_any:
                    _include_file(child)
                else:
                    suffix = os.path.splitext(child)[1].lower()
                    if suffix in {".txt", ".csv", ".xlsx", ".xls", ".docx", ".pptx", ".pdf", ".doc", ".ppt"}:
                        _include_file(child)
        if total_bytes > MAX_DIR_TOTAL_BYTES:
            raise RuntimeError(
                f"Directory exceeds total size limit ({MAX_DIR_TOTAL_BYTES // (1024 * 1024)} MB): {real_target}"
            )

    return files, warnings


def _build_csv_profile(path: str) -> dict:
    rows, fields = _read_csv_rows(path, limit=None)
    numeric_stats = {}
    for field in fields:
        numeric_values = []
        for row in rows:
            numeric_value = _coerce_float(row.get(field))
            if numeric_value is None:
                continue
            numeric_values.append(numeric_value)
        if numeric_values:
            numeric_stats[field] = {
                "count": len(numeric_values),
                "min": min(numeric_values),
                "max": max(numeric_values),
                "avg": round(sum(numeric_values) / len(numeric_values), 4),
            }
    numeric_fields = list(numeric_stats)
    return {
        "path": path,
        "fields": fields,
        "rowCountPreview": len(rows),
        "sampleRows": rows[:10],
        "numericStats": numeric_stats,
        "groupedNumericTotals": _build_grouped_numeric_totals(rows, fields, numeric_fields),
    }


def _build_workbook_profile(path: str) -> dict:
    workbook = _read_excel_workbook(path, limit_rows=120)
    sheet_profiles = {}
    for sheet_name, rows in workbook.items():
        headers = [str(value) for value in rows[0]] if rows else []
        records = []
        for row in rows[1:21]:
            record = {}
            for index, header in enumerate(headers):
                record[header] = row[index] if index < len(row) else None
            records.append(record)
        sheet_profiles[sheet_name] = {
            "headers": headers,
            "sampleRows": records,
            "rowCountPreview": max(0, len(rows) - 1),
        }
    return {"path": path, "sheets": sheet_profiles}


def _extract_txt_fragments(path: str, fragment_prefix: str | None = None) -> list[dict]:
    text = _read_txt(path)
    lines = text.splitlines()
    fragments = []
    current_title = ""
    current_lines: list[str] = []
    index = 0
    normalized_prefix = str(fragment_prefix or os.path.basename(path)).replace(os.sep, "/")

    def _flush():
        nonlocal index, current_title, current_lines
        body = "\n".join(current_lines).strip()
        if not current_title or not body:
            current_title = ""
            current_lines = []
            return
        index += 1
        fragments.append({
            "fragmentId": f"{normalized_prefix}::{index}",
            "title": current_title,
            "sourcePath": path,
            "content": body,
            "preview": body[:300],
        })
        current_title = ""
        current_lines = []

    for line in lines:
        if line.strip().startswith(">>>"):
            _flush()
            current_title = line.strip().lstrip("> ").strip()
            continue
        current_lines.append(line)
    _flush()
    return fragments


def _build_organize_context(target_paths: list[str]) -> tuple[dict, list[str]]:
    files, warnings = _collect_files(target_paths, allow_any=True)
    roots = [
        os.path.realpath(path if os.path.isdir(path) else os.path.dirname(path))
        for path in target_paths
    ]
    common_root = os.path.commonpath(roots)
    inventory = []
    fragments = []
    for path in files:
        relative = os.path.relpath(path, common_root)
        inventory.append({
            "sourcePath": path,
            "relativePath": relative,
            "extension": os.path.splitext(path)[1].lower(),
            "sizeBytes": os.path.getsize(path),
        })
        if path.lower().endswith(".txt"):
            fragments.extend(_extract_txt_fragments(path, fragment_prefix=relative))
    return {
        "commonRoot": common_root,
        "files": inventory,
        "fragments": fragments,
        "allowedActions": ["mkdir", "write_text", "write_fragment"],
    }, warnings


def _runtime_organize_context(organize_context: dict, preview_chars: int = 800) -> dict:
    runtime_fragments = []
    for fragment in organize_context.get("fragments", []):
        preview = str(fragment.get("preview") or fragment.get("content") or "")[:preview_chars]
        runtime_fragments.append({
            "fragmentId": fragment.get("fragmentId"),
            "title": fragment.get("title"),
            "sourcePath": fragment.get("sourcePath"),
            "preview": preview,
        })
    return {
        "commonRoot": organize_context.get("commonRoot"),
        "files": organize_context.get("files") or [],
        "fragments": runtime_fragments,
        "allowedActions": organize_context.get("allowedActions") or [],
    }


def _canonicalize_organize_destination(destination: str) -> str:
    relative = str(destination or "").strip().replace("\\", "/")
    if not relative or os.path.isabs(relative):
        raise RuntimeError(f"Unsafe organize destination: {destination}")
    parts = [part for part in relative.split("/") if part and part != "."]
    if not parts or ".." in parts:
        raise RuntimeError(f"Unsafe organize destination: {destination}")
    while len(parts) > 1 and parts[0].lower() in _ORGANIZE_WRAPPER_SEGMENTS:
        parts = parts[1:]
    normalized = "/".join(parts)
    if normalized == _ORGANIZE_SCHEMA_ROOT:
        raise RuntimeError(f"Unsafe organize destination: {destination}")
    if normalized.startswith(f"{_ORGANIZE_SCHEMA_ROOT}/"):
        return normalized
    return f"{_ORGANIZE_SCHEMA_ROOT}/{normalized}"


def _validate_actions(actions: list[dict], organize_context: dict) -> None:
    valid_fragments = {item["fragmentId"] for item in organize_context.get("fragments", [])}
    for action in actions:
        action_name = str(action.get("action") or "").strip()
        if action_name not in {"mkdir", "write_text", "write_fragment"}:
            raise RuntimeError(f"Unsupported organize action: {action_name}")
        _canonicalize_organize_destination(str(action.get("destination") or ""))
        if action_name == "write_fragment" and str(action.get("fragment_id") or "") not in valid_fragments:
            raise RuntimeError(f"Unknown organize fragment id: {action.get('fragment_id')}")


def _write_warnings(metadata: dict, warnings: list[str]) -> None:
    if warnings:
        _save_audit_file(metadata, "warnings.md", "\n".join(f"- {item}" for item in warnings) + "\n")


def _normalize_generated_text_content(content: str) -> str:
    text = str(content or "")
    if "\\n" not in text and "\\r" not in text and "\\t" not in text:
        return text
    if "\n" in text or "\r" in text:
        return text.replace("\\r\\n", "\n").replace("\\n", "\n").replace("\\t", "\t")
    return text.replace("\\r\\n", "\n").replace("\\n", "\n").replace("\\t", "\t")


def _non_overwrite_path(path: str) -> str:
    """Return path unchanged if it doesn't exist; otherwise append a compact timestamp suffix."""
    if not os.path.exists(path):
        return path
    base, ext = os.path.splitext(path)
    ts = local_iso_timestamp().replace(":", "").replace("-", "")[:15]
    candidate = f"{base}-{ts}{ext}"
    index = 1
    while os.path.exists(candidate):
        candidate = f"{base}-{ts}-{index}{ext}"
        index += 1
    return candidate


def _acquire_write_root(path: str) -> None:
    normalized = os.path.realpath(path)
    with _WRITE_ROOTS_LOCK:
        if normalized in _ACTIVE_WRITE_ROOTS:
            raise RuntimeError(f"Another Office task is already writing to: {normalized}")
        _ACTIVE_WRITE_ROOTS.add(normalized)


def _release_write_root(path: str) -> None:
    normalized = os.path.realpath(path)
    with _WRITE_ROOTS_LOCK:
        _ACTIVE_WRITE_ROOTS.discard(normalized)


def _artifact_metadata(metadata: dict, capability: str, task_id: str) -> dict:
    return {
        "agentId": AGENT_ID,
        "capability": capability,
        "taskId": task_id,
        "orchestratorTaskId": str(metadata.get("orchestratorTaskId") or task_id),
    }


def _target_output_dir(metadata: dict, target_paths: list[str], output_mode: str) -> str:
    if output_mode == "workspace":
        return _audit_dir(metadata)
    roots = [path if os.path.isdir(path) else os.path.dirname(path) for path in target_paths]
    output_root = os.path.commonpath(roots)
    os.makedirs(output_root, exist_ok=True)
    return output_root


def _write_manifest(output_root: str, task_id: str, output_mode: str, manifest_entries: list[dict]) -> None:
    manifest_path = os.path.join(output_root, ".office-agent-manifest.json")
    _write_text_file(
        manifest_path,
        json.dumps(
            {
                "taskId": task_id,
                "outputMode": output_mode,
                "executedActions": manifest_entries,
                "updatedAt": local_iso_timestamp(),
            },
            ensure_ascii=False,
            indent=2,
        ),
    )


def _execute_summary(capability: str, user_text: str, metadata: dict, task_id: str) -> dict:
    target_paths = list(metadata.get("officeTargetPaths") or [])
    # Preflight: check directory limits before reading content
    scan = _preflight_scan(target_paths)
    if scan["overFileCountLimit"] or scan["overBytesLimit"]:
        report = (
            f"Preflight limit exceeded — {scan['fileCount']} files, "
            f"{scan['totalBytes'] // (1024 * 1024)} MB total "
            f"(limits: {scan['limitFileCount']} files, {scan['limitTotalMB']} MB). "
            "Please narrow the target path or raise OFFICE_MAX_DIR_FILE_COUNT / OFFICE_MAX_DIR_TOTAL_MB."
        )
        return {
            "summary": report,
            "artifacts": [
                build_text_artifact(
                    "office-preflight-report",
                    report,
                    metadata=_artifact_metadata(metadata, capability, task_id),
                )
            ],
            "warnings": [],
        }

    files, warnings = _collect_files(target_paths)
    previews = []
    for path in files:
        try:
            previews.append(_extract_document_preview(path))
        except Exception as exc:
            warnings.append(f"Skipped {path}: {exc}")
    if not previews:
        raise RuntimeError("No readable files were available for summary.")

    response = _run_agentic_json(
        prompts.SUMMARIZE_TEMPLATE.format(user_text=user_text),
        "summarize",
        system_prompt=_build_manifest_prompt(__file__, prompts.SUMMARIZE_SYSTEM),
        context={"documents": previews},
    )
    summary_markdown = str(response.get("summary_markdown") or "").strip()
    if not summary_markdown:
        raise RuntimeError("Runtime did not return summary_markdown.")

    warnings.extend(str(item) for item in (response.get("warnings") or []) if str(item).strip())
    output_mode = str(metadata.get("officeOutputMode") or "workspace")
    output_root = _target_output_dir(metadata, target_paths, output_mode)
    base_filename = "summary.md"
    output_path = _non_overwrite_path(os.path.join(output_root, base_filename))
    _write_text_file(output_path, summary_markdown)
    _write_warnings(metadata, warnings)
    return {
        "summary": f"Summary created at {output_path}",
        "artifacts": [
            build_text_artifact(
                "office-summary",
                summary_markdown,
                metadata=_artifact_metadata(metadata, capability, task_id),
            )
        ],
        "warnings": warnings,
    }


def _execute_analysis(capability: str, user_text: str, metadata: dict, task_id: str) -> dict:
    target_paths = list(metadata.get("officeTargetPaths") or [])
    # Preflight: check directory limits before reading content
    scan = _preflight_scan(target_paths)
    if scan["overFileCountLimit"] or scan["overBytesLimit"]:
        report = (
            f"Preflight limit exceeded — {scan['fileCount']} files, "
            f"{scan['totalBytes'] // (1024 * 1024)} MB total "
            f"(limits: {scan['limitFileCount']} files, {scan['limitTotalMB']} MB)."
        )
        return {
            "summary": report,
            "artifacts": [
                build_text_artifact(
                    "office-preflight-report",
                    report,
                    metadata=_artifact_metadata(metadata, capability, task_id),
                )
            ],
            "warnings": [],
        }

    files, warnings = _collect_files(target_paths)
    profiles = []
    for path in files:
        suffix = os.path.splitext(path)[1].lower()
        try:
            if suffix == ".csv":
                profiles.append(_build_csv_profile(path))
            elif suffix in {".xlsx", ".xls"}:
                profiles.append(_build_workbook_profile(path))
        except Exception as exc:
            warnings.append(f"Skipped {path}: {exc}")
    if not profiles:
        raise RuntimeError("No readable CSV/XLSX/XLS files were available for analysis.")

    response = _run_agentic_json(
        prompts.ANALYZE_TEMPLATE.format(user_text=user_text),
        "analyze",
        system_prompt=_build_manifest_prompt(__file__, prompts.ANALYZE_SYSTEM),
        context={"datasets": profiles},
    )
    summary_markdown = str(response.get("summary_markdown") or "").strip()
    if not summary_markdown:
        raise RuntimeError("Runtime did not return summary_markdown.")

    warnings.extend(str(item) for item in (response.get("warnings") or []) if str(item).strip())
    output_mode = str(metadata.get("officeOutputMode") or "workspace")
    output_root = _target_output_dir(metadata, target_paths, output_mode)
    base_filename = "analysis.md"
    output_path = _non_overwrite_path(os.path.join(output_root, base_filename))
    _write_text_file(output_path, summary_markdown)
    _write_warnings(metadata, warnings)
    return {
        "summary": f"Analysis created at {output_path}",
        "artifacts": [
            build_text_artifact(
                "office-analysis",
                summary_markdown,
                metadata=_artifact_metadata(metadata, capability, task_id),
            )
        ],
        "warnings": warnings,
    }


def _execute_organize(capability: str, user_text: str, metadata: dict, task_id: str) -> dict:
    target_paths = list(metadata.get("officeTargetPaths") or [])
    organize_context, warnings = _build_organize_context(target_paths)
    response = _run_agentic_json(
        prompts.ORGANIZE_TEMPLATE.format(user_text=user_text),
        "organize",
        system_prompt=_build_manifest_prompt(__file__, prompts.ORGANIZE_SYSTEM),
        context={"inventory": _runtime_organize_context(organize_context)},
        timeout=300,
    )
    actions = response.get("actions") or []
    if not isinstance(actions, list):
        raise RuntimeError("Runtime did not return a valid actions list.")
    _validate_actions(actions, organize_context)

    output_mode = str(metadata.get("officeOutputMode") or "workspace")
    base_output_root = _target_output_dir(metadata, target_paths, output_mode)
    output_root = os.path.join(base_output_root, "organized-output")
    os.makedirs(output_root, exist_ok=True)

    # Persist the validated plan BEFORE any writes (required by design — R8 / §9.6).
    _save_audit_file(metadata, "operations-plan.json", json.dumps({"actions": actions}, ensure_ascii=False, indent=2))
    _append_command_log(metadata, f"Executing {len(actions)} organize action(s), outputMode={output_mode}")

    manifest_entries: list[dict] = []
    fragments = {item["fragmentId"]: item for item in organize_context.get("fragments", [])}
    for action in actions:
        action_name = action.get("action")
        requested_destination = _canonicalize_organize_destination(str(action.get("destination") or ""))
        destination = _safe_output_path(output_root, requested_destination)
        if action_name in {"write_text", "write_fragment"}:
            resolved_destination = _non_overwrite_path(destination)
            if resolved_destination != destination:
                warnings.append(
                    "Avoided overwrite for "
                    f"{requested_destination}; wrote to {os.path.relpath(resolved_destination, output_root)} instead."
                )
            destination = resolved_destination
        if action_name == "mkdir":
            os.makedirs(destination, exist_ok=True)
        elif action_name == "write_text":
            _write_text_file(destination, _normalize_generated_text_content(str(action.get("content") or "")))
        elif action_name == "write_fragment":
            fragment = fragments[str(action.get("fragment_id"))]
            _write_text_file(destination, fragment.get("content", ""))
        manifest_entries.append({
            "action": action_name,
            "destination": destination,
            "relativeDestination": requested_destination,
            "ts": local_iso_timestamp(),
        })
        _write_manifest(output_root, task_id, output_mode, manifest_entries)
        # For inplace mode append per-step progress to command-log for human recoverability (§9.6).
        if output_mode == "inplace":
            _append_command_log(metadata, f"Executed {action_name} → {destination}")

    summary_markdown = str(response.get("summary_markdown") or "").strip()
    if not summary_markdown:
        summary_markdown = f"Organize plan executed with {len(actions)} action(s)."
    warnings.extend(str(item) for item in (response.get("warnings") or []) if str(item).strip())
    _write_text_file(os.path.join(_audit_dir(metadata), "organize-report.md"), summary_markdown)
    _write_warnings(metadata, warnings)
    return {
        "summary": f"Organize plan executed at {output_root}",
        "artifacts": [
            build_text_artifact(
                "office-organize-report",
                summary_markdown,
                metadata=_artifact_metadata(metadata, capability, task_id),
            )
        ],
        "warnings": warnings,
    }


def _execute_capability(task_id: str, message: dict) -> dict:
    metadata = dict(message.get("metadata") or {})
    metadata.setdefault("taskId", task_id)
    capability = str(metadata.get("requestedCapability") or "").strip()
    user_text = extract_text(message)
    input_root = str(metadata.get("officeInputRoot") or INPUT_ROOT).strip() or INPUT_ROOT
    output_mode = str(metadata.get("officeOutputMode") or "workspace")
    if not capability:
        raise RuntimeError("Office task is missing requestedCapability metadata.")

    target_paths = [os.path.realpath(path) for path in (metadata.get("officeTargetPaths") or []) if str(path).strip()]
    if not target_paths:
        raise RuntimeError("Office task is missing officeTargetPaths metadata.")

    metadata["officeTargetPaths"] = target_paths
    _enforce_office_task_permissions(metadata, target_paths, input_root, output_mode)
    write_root = ""
    if output_mode == "inplace":
        write_root = _target_output_dir(metadata, target_paths, output_mode)
        _acquire_write_root(write_root)

    try:
        _append_command_log(metadata, f"Starting capability {capability}")
        _update_stage_summary(metadata, task_id, "preflight", capability=capability, outputMode=metadata.get("officeOutputMode"))

        if capability in {"office.document.summarize", "office.folder.summarize"}:
            _update_stage_summary(metadata, task_id, "summarizing", capability=capability)
            return _execute_summary(capability, user_text, metadata, task_id)
        if capability == "office.data.analyze":
            _update_stage_summary(metadata, task_id, "analyzing", capability=capability)
            return _execute_analysis(capability, user_text, metadata, task_id)
        if capability == "office.folder.organize":
            _update_stage_summary(metadata, task_id, "organizing", capability=capability)
            return _execute_organize(capability, user_text, metadata, task_id)
        raise RuntimeError(f"Unsupported Office capability: {capability}")
    finally:
        if write_root:
            _release_write_root(write_root)


def _build_office_task_prompt(
    *,
    user_text: str,
    capability: str,
    target_paths: list[str],
    output_mode: str,
    workspace_path: str,
    task_id: str,
    compass_task_id: str,
) -> str:
    """Build the task prompt for the agentic runtime."""
    paths_text = "\n".join(f"- {p}" for p in target_paths) if target_paths else "- (not specified)"
    output_section = ""
    if output_mode == "inplace":
        output_section = "Output mode: IN-PLACE — write results back into the source directory."
    elif workspace_path:
        output_section = f"Output mode: WORKSPACE — write results to {workspace_path}/office-agent/"
    else:
        output_section = "Output mode: return summary in task result."

    task_description = ""
    if capability in {"office.document.summarize", "office.folder.summarize"}:
        task_description = """Task: DOCUMENT SUMMARIZATION
- Use read_file to read each target document
- For each document: identify title, key topics, main points, recommendations
- Generate a concise summary (max 500 words per document)
- If multiple documents: also create an overall cross-document summary
- Save summary to output location as summary.md"""
    elif capability == "office.data.analyze":
        task_description = """Task: DATA ANALYSIS
- Use read_file or bash (cat/head) to read CSV/spreadsheet files
- Identify columns, data types, value ranges, missing data
- Compute basic statistics (counts, averages, distributions)
- Identify patterns, anomalies, correlations
- Generate an analysis report with key findings
- Save report to output location as analysis.md"""
    elif capability == "office.folder.organize":
        task_description = """Task: FOLDER ORGANIZATION
- Use bash (ls -la, find) to enumerate files in target directories
- Group files by type, topic, date, or logical category
- Use write_file to create a reorganization plan (plan.json)
- For inplace mode: move files to organized subdirectories using bash
- For workspace mode: write the organization plan only, do not move files
- Save organization report to output location as organization-report.md"""
    else:
        task_description = f"Task: {capability}\n{user_text}"

    return f"""You are the Office Agent. Process the following office/document task.

## User Request
{user_text}

## Capability
{capability}

## Target Files/Directories
{paths_text}

## {task_description}

## Output
{output_section}
- orchestratorTaskId: {compass_task_id}
- officeAgentTaskId: {task_id}

## Rules
- Only access files within the specified target paths
- Do NOT access files outside the target paths
- For summarize/analyze: read files, do not modify source files unless output_mode=inplace
- Use read_file for text files, use bash (cat) for binary inspection
- After completing work, use complete_current_task with a summary and any output file paths as artifacts
- If you cannot access a file (permission denied, binary, too large), note it in warnings
- Use report_progress to announce: "Reading files", "Analyzing content", "Writing results"
"""


def _run_workflow(task_id: str, message: dict):
    metadata = dict(message.get("metadata") or {})
    callback_url = str(metadata.get("orchestratorCallbackUrl") or "")
    compass_url = resolve_orchestrator_base_url(metadata)
    compass_task_id = str(metadata.get("orchestratorTaskId") or "")
    exit_rule = PerTaskExitHandler.parse(metadata)
    workspace_path = str(metadata.get("sharedWorkspacePath") or "")
    task = task_store.get(task_id)
    if not task:
        return

    capability = str(metadata.get("requestedCapability") or "").strip()
    user_text = extract_text(message)
    input_root = str(metadata.get("officeInputRoot") or INPUT_ROOT).strip() or INPUT_ROOT
    output_mode = str(metadata.get("officeOutputMode") or "workspace")
    target_paths = [os.path.realpath(p) for p in (metadata.get("officeTargetPaths") or []) if str(p).strip()]

    # Validate permissions before starting agentic execution
    if target_paths:
        try:
            _enforce_office_task_permissions(metadata, target_paths, input_root, output_mode)
        except (PermissionDeniedError, RuntimeError) as perm_err:
            failure = f"Office task failed: {perm_err}"
            artifacts = []
            if isinstance(perm_err, PermissionDeniedError):
                artifacts = [build_permission_denied_artifact(perm_err.details, agent_id=AGENT_ID)]
            task_store.update_state(task_id, "TASK_STATE_FAILED", failure)
            task.artifacts = artifacts
            _notify_callback(callback_url, task_id, "TASK_STATE_FAILED", failure, artifacts)
            audit_log("TASK_FAILED", task_id=task_id, error=str(perm_err))
            _apply_task_exit_rule(task_id, exit_rule)
            return

    configure_control_tools(
        task_context={
            "taskId": task_id,
            "agentId": AGENT_ID,
            "workspacePath": workspace_path,
            "permissions": metadata.get("permissions"),
        },
        complete_fn=lambda result, artifacts: task_store.update_state(task_id, "TASK_STATE_COMPLETED", result),
        fail_fn=lambda error: task_store.update_state(task_id, "TASK_STATE_FAILED", error),
        input_required_fn=lambda question, ctx: task_store.update_state(task_id, "TASK_STATE_INPUT_REQUIRED", question),
    )

    task_store.update_state(task_id, "TASK_STATE_WORKING", "Office Agent is processing the request.")
    _report_progress(compass_url, compass_task_id, "Office Agent is starting.")
    audit_log("TASK_STARTED", task_id=task_id, capability=capability)

    # Build system prompt from manifest or rules
    manifest_prompt = _get_manifest_system_prompt()
    system_prompt = manifest_prompt or build_system_prompt(
        prompts.SYSTEM,
        "office",
    )

    task_prompt = _build_office_task_prompt(
        user_text=user_text,
        capability=capability,
        target_paths=target_paths,
        output_mode=output_mode,
        workspace_path=workspace_path,
        task_id=task_id,
        compass_task_id=compass_task_id,
    )

    # Set sandbox root to allow reading from target paths and writing to workspace
    cwd = target_paths[0] if target_paths else workspace_path or os.getcwd()

    try:
        runtime = get_runtime()
        result = runtime.run_agentic(
            task=task_prompt,
            system_prompt=system_prompt,
            cwd=cwd,
            max_turns=40,
            timeout=TASK_TIMEOUT,
        )
        summary = result.summary or "Office task completed."
        final_artifacts: list = list(result.artifacts or [])

        summary_artifact = {
            "name": "office-agent-summary",
            "artifactType": "text/plain",
            "parts": [{"text": summary}],
            "metadata": {
                "agentId": AGENT_ID,
                "capability": capability,
                "orchestratorTaskId": compass_task_id,
                "taskId": task_id,
            },
        }
        final_artifacts.insert(0, summary_artifact)

        if result.success:
            task_store.update_state(task_id, "TASK_STATE_COMPLETED", summary)
            task.artifacts = final_artifacts
            _update_stage_summary(metadata, task_id, "completed")
            _report_progress(compass_url, compass_task_id, summary)
            _notify_callback(callback_url, task_id, "TASK_STATE_COMPLETED", summary, final_artifacts)
            audit_log("TASK_COMPLETED", task_id=task_id, capability=capability)
        else:
            task_store.update_state(task_id, "TASK_STATE_FAILED", summary)
            task.artifacts = final_artifacts
            _update_stage_summary(metadata, task_id, "failed", error=summary)
            _report_progress(compass_url, compass_task_id, summary)
            _notify_callback(callback_url, task_id, "TASK_STATE_FAILED", summary, final_artifacts)
            audit_log("TASK_FAILED", task_id=task_id, error=summary[:300])

    except Exception as err:
        failure = f"Office task failed: {err}"
        artifacts = []
        if isinstance(err, PermissionDeniedError):
            artifacts = [build_permission_denied_artifact(err.details, agent_id=AGENT_ID)]
        task_store.update_state(task_id, "TASK_STATE_FAILED", failure)
        task.artifacts = artifacts
        _save_audit_file(metadata, "failure.txt", failure + "\n")
        _update_stage_summary(metadata, task_id, "failed", error=str(err))
        _report_progress(compass_url, compass_task_id, failure)
        _notify_callback(callback_url, task_id, "TASK_STATE_FAILED", failure, artifacts)
        audit_log("TASK_FAILED", task_id=task_id, error=str(err))
    finally:
        _apply_task_exit_rule(task_id, exit_rule)


class OfficeAgentHandler(BaseHTTPRequestHandler):
    protocol_version = "HTTP/1.1"

    def _send_json(self, code: int, payload: dict):
        body = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self.send_response(code)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _read_body(self) -> dict:
        length = int(self.headers.get("Content-Length", 0))
        if length <= 0:
            return {}
        return json.loads(self.rfile.read(length).decode("utf-8"))

    def do_GET(self):
        path = urlparse(self.path).path
        if path == "/health":
            self._send_json(200, {"status": "ok", "service": AGENT_ID})
            return
        if path == "/.well-known/agent-card.json":
            card_path = os.path.join(os.path.dirname(__file__), "agent-card.json")
            with open(card_path, encoding="utf-8") as handle:
                card = json.load(handle)
            text = json.dumps(card).replace("__ADVERTISED_URL__", ADVERTISED_URL)
            self._send_json(200, json.loads(text))
            return
        match = re.fullmatch(r"/tasks/([^/]+)", path)
        if match:
            task = task_store.get(match.group(1))
            if not task:
                self._send_json(404, {"error": "task_not_found"})
                return
            self._send_json(200, {"task": task.to_dict()})
            return
        self._send_json(404, {"error": "not_found"})

    def do_POST(self):
        path = urlparse(self.path).path
        match = re.fullmatch(r"/tasks/([^/]+)/ack", path)
        if match:
            task_id = match.group(1)
            acked = exit_handler.acknowledge(task_id)
            print(f"[{AGENT_ID}] Received ACK for task {task_id} (registered={acked})")
            self._send_json(200, {"ok": True, "task_id": task_id})
            return

        if path != "/message:send":
            self._send_json(404, {"error": "not_found"})
            return

        body = self._read_body()
        message = body.get("message") or {}
        if not message:
            self._send_json(400, {"error": "missing_message"})
            return

        task = task_store.create()
        audit_log(
            "TASK_RECEIVED",
            task_id=task.task_id,
            capability=(message.get("metadata") or {}).get("requestedCapability", ""),
            instruction_preview=extract_text(message)[:120],
        )
        worker = threading.Thread(target=_run_workflow, args=(task.task_id, message), daemon=True)
        worker.start()
        self._send_json(200, {"task": task.to_dict()})

    def log_message(self, fmt, *args):
        line = args[0] if args else ""
        if any(part in line for part in ("/health", "/.well-known/agent-card.json")):
            return
        print(
            f"[{AGENT_ID}] {line} "
            f"{args[1] if len(args) > 1 else ''} "
            f"{args[2] if len(args) > 2 else ''}"
        )


def main():
    global _SERVER
    print(f"[{AGENT_ID}] Office Agent starting on {HOST}:{PORT}")
    reporter.start()
    _SERVER = ThreadingHTTPServer((HOST, PORT), OfficeAgentHandler)
    _SERVER.serve_forever()


if __name__ == "__main__":
    main()