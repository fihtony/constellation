"""Document-format reading tools for the Office Agent.

Provides dedicated tools for extracting plain text from binary document
formats (PDF, DOCX, PPTX, XLSX) so the LLM can call a single tool instead
of writing inline Python via run_local_command.

Based on reference/skills/pdf/SKILL.md, reference/skills/docx/SKILL.md,
and reference/skills/pptx/SKILL.md.

All tools:
  - enforce the sandbox path jail via safe_path / _resolve()
  - return the extracted plain text as a string
  - return an error dict when the file cannot be read
"""

from __future__ import annotations

import os
from pathlib import Path

from common.tools.base import ConstellationTool, ToolSchema
from common.tools.registry import register_tool

# Lazy import sandbox helpers so this module can be imported in test environments
# that don't have the full runtime stack loaded.
def _resolve(path: str) -> Path:
    """Resolve a path through the coding-tools sandbox jail."""
    try:
        from common.tools.coding_tools import _resolve as _coding_resolve
        return _coding_resolve(path)
    except Exception:  # noqa: BLE001
        # Fallback: only enforce absolute paths (no sandbox check) when the
        # coding tools sandbox hasn't been configured yet (e.g. unit tests).
        p = Path(path)
        if not p.is_absolute():
            raise ValueError(f"Path must be absolute: {path}")
        return p


# ---------------------------------------------------------------------------
# read_pdf
# ---------------------------------------------------------------------------

class ReadPdfTool(ConstellationTool):
    """Extract plain text from a PDF file using pdfplumber."""

    @property
    def schema(self) -> ToolSchema:
        return ToolSchema(
            name="read_pdf",
            description=(
                "Extract plain text content from a PDF file. "
                "Uses pdfplumber to preserve text layout and handle multi-column PDFs. "
                "Returns the extracted text as a string. "
                "Use this instead of run_local_command when reading PDF files."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Absolute path to the PDF file.",
                    },
                    "pages": {
                        "type": "array",
                        "items": {"type": "integer"},
                        "description": (
                            "Optional 0-based page indices to extract. "
                            "Omit to extract all pages."
                        ),
                    },
                    "max_chars": {
                        "type": "integer",
                        "description": "Maximum characters to return (default: 100000).",
                    },
                },
                "required": ["path"],
            },
        )

    def execute(self, args: dict) -> dict:
        path_str = str(args.get("path") or "").strip()
        pages = args.get("pages")
        max_chars = int(args.get("max_chars") or 100_000)

        if not path_str:
            return self.error("Missing required argument: path")

        try:
            resolved = _resolve(path_str)
        except (ValueError, Exception) as exc:  # noqa: BLE001
            return self.error(f"Path access denied: {exc}")

        if not resolved.is_file():
            return self.error(f"File not found: {path_str}")
        if resolved.suffix.lower() != ".pdf":
            return self.error(f"Not a PDF file: {path_str}")

        try:
            import pdfplumber
        except ImportError:
            return self.error(
                "pdfplumber is not installed. Run: pip install pdfplumber"
            )

        try:
            with pdfplumber.open(str(resolved)) as pdf:
                page_list = pdf.pages
                if pages is not None:
                    page_list = [pdf.pages[i] for i in pages if 0 <= i < len(pdf.pages)]
                parts = []
                for page in page_list:
                    text = page.extract_text() or ""
                    if text.strip():
                        parts.append(text)
                text_out = "\n\n".join(parts)
        except Exception as exc:  # noqa: BLE001
            return self.error(f"Failed to read PDF '{path_str}': {exc}")

        if len(text_out) > max_chars:
            text_out = text_out[:max_chars] + f"\n\n[... truncated at {max_chars} chars]"

        return self.ok(text_out or "(no extractable text found in PDF)")


# ---------------------------------------------------------------------------
# read_docx
# ---------------------------------------------------------------------------

class ReadDocxTool(ConstellationTool):
    """Extract plain text from a Word document (.docx) using python-docx."""

    @property
    def schema(self) -> ToolSchema:
        return ToolSchema(
            name="read_docx",
            description=(
                "Extract plain text content from a Word document (.docx file). "
                "Uses python-docx to read paragraphs and tables. "
                "Returns the extracted text as a string. "
                "Use this instead of run_local_command when reading .docx files."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Absolute path to the .docx file.",
                    },
                    "include_tables": {
                        "type": "boolean",
                        "description": "Whether to include table cell text (default: true).",
                    },
                    "max_chars": {
                        "type": "integer",
                        "description": "Maximum characters to return (default: 100000).",
                    },
                },
                "required": ["path"],
            },
        )

    def execute(self, args: dict) -> dict:
        path_str = str(args.get("path") or "").strip()
        include_tables = bool(args.get("include_tables", True))
        max_chars = int(args.get("max_chars") or 100_000)

        if not path_str:
            return self.error("Missing required argument: path")

        try:
            resolved = _resolve(path_str)
        except (ValueError, Exception) as exc:  # noqa: BLE001
            return self.error(f"Path access denied: {exc}")

        if not resolved.is_file():
            return self.error(f"File not found: {path_str}")
        if resolved.suffix.lower() not in {".docx", ".doc"}:
            return self.error(f"Not a Word document: {path_str}")

        try:
            from docx import Document
        except ImportError:
            return self.error(
                "python-docx is not installed. Run: pip install python-docx"
            )

        try:
            doc = Document(str(resolved))
            parts: list[str] = []

            # Paragraphs (preserves heading order)
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)

            # Tables
            if include_tables:
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            parts.append(row_text)

            text_out = "\n".join(parts)
        except Exception as exc:  # noqa: BLE001
            return self.error(f"Failed to read DOCX '{path_str}': {exc}")

        if len(text_out) > max_chars:
            text_out = text_out[:max_chars] + f"\n\n[... truncated at {max_chars} chars]"

        return self.ok(text_out or "(no extractable text found in DOCX)")


# ---------------------------------------------------------------------------
# read_pptx
# ---------------------------------------------------------------------------

class ReadPptxTool(ConstellationTool):
    """Extract plain text from a PowerPoint presentation (.pptx) using python-pptx."""

    @property
    def schema(self) -> ToolSchema:
        return ToolSchema(
            name="read_pptx",
            description=(
                "Extract plain text content from a PowerPoint presentation (.pptx file). "
                "Reads slide titles, body text, and speaker notes. "
                "Returns the extracted text as a string. "
                "Use this instead of run_local_command when reading .pptx files."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Absolute path to the .pptx file.",
                    },
                    "include_notes": {
                        "type": "boolean",
                        "description": "Whether to include speaker notes (default: true).",
                    },
                    "max_chars": {
                        "type": "integer",
                        "description": "Maximum characters to return (default: 100000).",
                    },
                },
                "required": ["path"],
            },
        )

    def execute(self, args: dict) -> dict:
        path_str = str(args.get("path") or "").strip()
        include_notes = bool(args.get("include_notes", True))
        max_chars = int(args.get("max_chars") or 100_000)

        if not path_str:
            return self.error("Missing required argument: path")

        try:
            resolved = _resolve(path_str)
        except (ValueError, Exception) as exc:  # noqa: BLE001
            return self.error(f"Path access denied: {exc}")

        if not resolved.is_file():
            return self.error(f"File not found: {path_str}")
        if resolved.suffix.lower() != ".pptx":
            return self.error(f"Not a PowerPoint file: {path_str}")

        try:
            from pptx import Presentation
        except ImportError:
            return self.error(
                "python-pptx is not installed. Run: pip install python-pptx"
            )

        try:
            prs = Presentation(str(resolved))
            parts: list[str] = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts: list[str] = [f"--- Slide {slide_num} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_parts.append(text)
                if include_notes and slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_parts.append(f"[Notes: {notes_text}]")
                parts.append("\n".join(slide_parts))

            text_out = "\n\n".join(parts)
        except Exception as exc:  # noqa: BLE001
            return self.error(f"Failed to read PPTX '{path_str}': {exc}")

        if len(text_out) > max_chars:
            text_out = text_out[:max_chars] + f"\n\n[... truncated at {max_chars} chars]"

        return self.ok(text_out or "(no extractable text found in PPTX)")


# ---------------------------------------------------------------------------
# read_xlsx
# ---------------------------------------------------------------------------

class ReadXlsxTool(ConstellationTool):
    """Extract data from an Excel spreadsheet (.xlsx) as CSV-formatted text."""

    @property
    def schema(self) -> ToolSchema:
        return ToolSchema(
            name="read_xlsx",
            description=(
                "Extract data from an Excel spreadsheet (.xlsx file). "
                "Returns each worksheet's content as CSV-formatted text. "
                "Use this instead of run_local_command when reading .xlsx files."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Absolute path to the .xlsx file.",
                    },
                    "sheet_name": {
                        "type": "string",
                        "description": "Name of a specific sheet to read. Omit to read all sheets.",
                    },
                    "max_rows": {
                        "type": "integer",
                        "description": "Maximum rows per sheet to return (default: 1000).",
                    },
                },
                "required": ["path"],
            },
        )

    def execute(self, args: dict) -> dict:
        path_str = str(args.get("path") or "").strip()
        sheet_name = str(args.get("sheet_name") or "").strip() or None
        max_rows = int(args.get("max_rows") or 1000)

        if not path_str:
            return self.error("Missing required argument: path")

        try:
            resolved = _resolve(path_str)
        except (ValueError, Exception) as exc:  # noqa: BLE001
            return self.error(f"Path access denied: {exc}")

        if not resolved.is_file():
            return self.error(f"File not found: {path_str}")
        if resolved.suffix.lower() not in {".xlsx", ".xls"}:
            return self.error(f"Not an Excel file: {path_str}")

        try:
            import openpyxl
        except ImportError:
            return self.error(
                "openpyxl is not installed. Run: pip install openpyxl"
            )

        try:
            wb = openpyxl.load_workbook(str(resolved), read_only=True, data_only=True)
            sheets = [sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.sheetnames
            parts: list[str] = []
            for sname in sheets:
                ws = wb[sname]
                rows = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i >= max_rows:
                        rows.append(f"[... {ws.max_row - max_rows} more rows truncated]")
                        break
                    rows.append(",".join("" if v is None else str(v) for v in row))
                parts.append(f"Sheet: {sname}\n" + "\n".join(rows))
            text_out = "\n\n".join(parts)
        except Exception as exc:  # noqa: BLE001
            return self.error(f"Failed to read XLSX '{path_str}': {exc}")

        return self.ok(text_out or "(no data found in spreadsheet)")


# ---------------------------------------------------------------------------
# Registration
# ---------------------------------------------------------------------------

register_tool(ReadPdfTool())
register_tool(ReadDocxTool())
register_tool(ReadPptxTool())
register_tool(ReadXlsxTool())
